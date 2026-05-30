import gradio as gr
import os
import re
import hashlib
import sqlite3
import json
import uuid
import shutil
import tempfile
import zipfile
from pathlib import Path
from datetime import datetime
from typing import Generator, Optional

# ─── Core ML / RAG ─────────────────────────────────────────────────────────────
from huggingface_hub import InferenceClient
from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, TextLoader, CSVLoader, UnstructuredMarkdownLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document

try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import pytesseract
    from pdf2image import convert_from_path
    from PIL import Image, ImageDraw, ImageFont
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
# Pillow alone (needed for slide images — no OCR required)
try:
    from PIL import Image, ImageDraw, ImageFont
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False

# Full OCR stack (pytesseract + pdf2image + Pillow)
try:
    import pytesseract
    from pdf2image import convert_from_path
    OCR_AVAILABLE = PILLOW_AVAILABLE
except ImportError:
    OCR_AVAILABLE = False
try:
    from rank_bm25 import BM25Okapi
    BM25_AVAILABLE = True
except ImportError:
    BM25_AVAILABLE = False

try:
    from sentence_transformers import CrossEncoder
    RERANKER_AVAILABLE = True
except ImportError:
    RERANKER_AVAILABLE = False

try:
    from fastapi import FastAPI, UploadFile, File, Body, Form, Header, HTTPException, Depends
    FASTAPI_AVAILABLE = True
except Exception:
    FASTAPI_AVAILABLE = False

# ─── Audio / Video ─────────────────────────────────────────────────────────────
try:
    import azure.cognitiveservices.speech as speechsdk
    AZURE_TTS_AVAILABLE = True
except Exception:
    AZURE_TTS_AVAILABLE = False

try:
    from gtts import gTTS
    GTTS_AVAILABLE = True
except Exception:
    GTTS_AVAILABLE = False

MOVIEPY_AVAILABLE = False
try:
    from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
    MOVIEPY_AVAILABLE = True
except Exception:
    try:
        import subprocess
        subprocess.run(
            ["pip", "install", "moviepy==1.0.3", "decorator<5.0", "--quiet", "--no-warn-script-location"],
            check=True, capture_output=True
        )
        from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
        MOVIEPY_AVAILABLE = True
    except Exception as _e:
        MOVIEPY_AVAILABLE = False
 
# ==============================================================================
# ⚙️  CONFIGURATION
# ==============================================================================

REPO_ID      = "meta-llama/Llama-3.3-70B-Instruct"
EMBED_MODEL  = "sentence-transformers/all-MiniLM-L6-v2"
RERANK_MODEL = "cross-encoder/ms-marco-MiniLM-L-6-v2"

DATA_DIR     = Path("/data")
DB_PATH      = str(DATA_DIR / "phoenix_history.db")
FALLBACK_SESSION_DIR = "shared"

MAX_HISTORY_TURNS = 12
MAX_CONTEXT_CHARS = 6500
SUMMARY_EVERY_N_ASSISTANT_TURNS = 6

HF_TOKEN = os.getenv("HF_TOKEN")
if not HF_TOKEN:
    raise ValueError("❌ HF_TOKEN missing!")

# The key/region go as the DEFAULT, not as the var name
AZURE_SPEECH_KEY    = os.getenv("AZURE_SPEECH_KEY",    "??")
AZURE_SPEECH_REGION = os.getenv("AZURE_SPEECH_REGION", "francecentral")
AI_SERVICE_TOKEN    = os.getenv("AI_SERVICE_TOKEN", "")
if not AI_SERVICE_TOKEN:
    print("⚠️ AI_SERVICE_TOKEN not set; internal API auth is disabled.")

client = InferenceClient(model=REPO_ID, token=HF_TOKEN)

_embeddings = None
def get_embeddings():
    global _embeddings
    if _embeddings is None:
        _embeddings = HuggingFaceEmbeddings(model_name=EMBED_MODEL)
    return _embeddings

_reranker = None
def get_reranker():
    global _reranker
    if _reranker is None and RERANKER_AVAILABLE:
        try:
            _reranker = CrossEncoder(RERANK_MODEL)
        except Exception:
            pass
    return _reranker

USER_CONTEXTS: dict = {}
USER_ACTIVE_SESSIONS: dict = {}

def get_or_create_session_id(username: str, session_id: Optional[str] = None) -> str:
    uname = sanitize_username(username)
    if session_id:
        return ensure_session(session_id, name=f"{uname or 'Session'}")
    if uname not in USER_ACTIVE_SESSIONS:
        USER_ACTIVE_SESSIONS[uname] = create_session(name=f"{uname or 'Session'}")
    return USER_ACTIVE_SESSIONS[uname]

def _context_key(username: str, session_id: Optional[str]) -> str:
    uname = sanitize_username(username)
    sid = get_or_create_session_id(username, session_id)
    return f"{uname}:{sid}"

def get_user_context(username: str, session_id: Optional[str] = None) -> dict:
    key = _context_key(username, session_id)
    if key not in USER_CONTEXTS:
        USER_CONTEXTS[key] = {"doc_store": {}, "active_docs": [], "doc_hashes": set()}
    return USER_CONTEXTS[key]

def clear_user_context(username: str, session_id: Optional[str] = None):
    key = _context_key(username, session_id)
    USER_CONTEXTS[key] = {"doc_store": {}, "active_docs": [], "doc_hashes": set()}

# ==============================================================================
# ✅ USERNAME + DATA DIRS
# ==============================================================================

def sanitize_username(name: str) -> str:
    name = (name or "").strip()
    name = re.sub(r"[^a-zA-Z0-9_\-]", "_", name)
    return name[:32] if name else ""

def is_valid_username(name: str) -> bool:
    cleaned = sanitize_username(name)
    return len(cleaned) >= 2

def ensure_user_dirs(username: str, session_id: Optional[str] = None) -> Path:
    DATA_DIR.mkdir(exist_ok=True)
    user_root = DATA_DIR / sanitize_username(username)
    (user_root / "history").mkdir(parents=True, exist_ok=True)
    uploads_root = user_root / "uploads"
    uploads_root.mkdir(parents=True, exist_ok=True)
    if session_id:
        (uploads_root / session_id).mkdir(parents=True, exist_ok=True)
    return user_root

# ==============================================================================
# 🗄️  DATABASE
# ==============================================================================

def init_db():
    DATA_DIR.mkdir(exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    conn.execute("PRAGMA journal_mode=WAL;")
    conn.execute("PRAGMA synchronous=NORMAL;")
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS sessions (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            doc_names TEXT DEFAULT '[]',
            summary TEXT DEFAULT ''
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            sources TEXT DEFAULT '[]',
            confidence REAL DEFAULT 0.0,
            created_at TEXT NOT NULL,
            FOREIGN KEY (session_id) REFERENCES sessions(id)
        )
    """)
    c.execute("CREATE INDEX IF NOT EXISTS idx_messages_session ON messages(session_id)")
    c.execute("CREATE INDEX IF NOT EXISTS idx_messages_content ON messages(content)")
    conn.commit()
    conn.close()

def create_session(name: str = None, doc_names: list = None, session_id: Optional[str] = None) -> str:
    session_id = session_id or str(uuid.uuid4())[:8]
    now = datetime.utcnow().isoformat()
    if not name:
        name = f"Session {datetime.now().strftime('%b %d %H:%M')}"
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "INSERT INTO sessions (id, name, created_at, updated_at, doc_names) VALUES (?,?,?,?,?)",
        (session_id, name, now, now, json.dumps(doc_names or []))
    )
    conn.commit()
    conn.close()
    return session_id

def ensure_session(session_id: str, name: str = None, doc_names: list = None) -> str:
    if not session_id:
        raise ValueError("session_id is required")
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute("SELECT id FROM sessions WHERE id=?", (session_id,)).fetchone()
    if not row:
        now = datetime.utcnow().isoformat()
        if not name:
            name = f"Session {datetime.now().strftime('%b %d %H:%M')}"
        conn.execute(
            "INSERT INTO sessions (id, name, created_at, updated_at, doc_names) VALUES (?,?,?,?,?)",
            (session_id, name, now, now, json.dumps(doc_names or []))
        )
        conn.commit()
    conn.close()
    return session_id

def save_message(session_id: str, role: str, content: str,
                 sources: list = None, confidence: float = 0.0):
    now = datetime.utcnow().isoformat()
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "INSERT INTO messages (session_id, role, content, sources, confidence, created_at) VALUES (?,?,?,?,?,?)",
        (session_id, role, content, json.dumps(sources or []), confidence, now)
    )
    conn.execute("UPDATE sessions SET updated_at=? WHERE id=?", (now, session_id))
    conn.commit()
    conn.close()

def get_session_summary(session_id: str) -> str:
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute("SELECT summary FROM sessions WHERE id=?", (session_id,)).fetchone()
    conn.close()
    return row[0] if row else ""

def update_session_summary(session_id: str, summary: str):
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "UPDATE sessions SET summary=?, updated_at=? WHERE id=?",
        (summary, datetime.utcnow().isoformat(), session_id)
    )
    conn.commit()
    conn.close()

def load_session_messages(session_id: str, limit: int = 40) -> list:
    conn = sqlite3.connect(DB_PATH)
    rows = conn.execute(
        "SELECT role, content, sources, confidence, created_at FROM messages "
        "WHERE session_id=? ORDER BY id DESC LIMIT ?",
        (session_id, limit)
    ).fetchall()
    conn.close()
    rows.reverse()
    return [
        {"role": r[0], "content": r[1],
         "sources": json.loads(r[2]), "confidence": r[3], "created_at": r[4]}
        for r in rows
    ]

def count_assistant_turns(session_id: str) -> int:
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute(
        "SELECT COUNT(*) FROM messages WHERE session_id=? AND role='assistant'",
        (session_id,)
    ).fetchone()
    conn.close()
    return row[0] if row else 0

def update_session_docs(session_id: str, doc_names: list):
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        "UPDATE sessions SET doc_names=?, updated_at=? WHERE id=?",
        (json.dumps(doc_names), datetime.utcnow().isoformat(), session_id)
    )
    conn.commit()
    conn.close()

init_db()

# ==============================================================================
# 🧾  USER HISTORY FILES (JSONL)
# ==============================================================================

def append_history_file(username: str, session_id: str, role: str, content: str):
    user_root = ensure_user_dirs(username)
    path = user_root / "history" / f"{session_id}.jsonl"
    entry = {"ts": datetime.utcnow().isoformat(), "role": role, "content": content}
    with open(path, "a", encoding="utf-8") as f:
        f.write(json.dumps(entry, ensure_ascii=False) + "\n")

def make_history_zip_for_user(username: str) -> Optional[str]:
    if sanitize_username(username) != "187177":
        return None
    user_root = ensure_user_dirs(username)
    hist_dir = user_root / "history"
    if not hist_dir.exists():
        return None
    zip_path = user_root / "history_export.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in hist_dir.glob("*.jsonl"):
            zf.write(p, arcname=p.name)
    return str(zip_path)

def get_db_path_for_download() -> str:
    return DB_PATH

# ==============================================================================
# 🔧  FILE LOADERS
# ==============================================================================

def load_pdf(file_path):
    docs = []
    if PDFPLUMBER_AVAILABLE:
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    for table in page.extract_tables():
                        rows = [" | ".join(str(c) for c in row if c) for row in table if row]
                        text += "\n" + "\n".join(rows)
                    if text.strip():
                        docs.append(Document(page_content=text,
                                             metadata={"page": i+1, "source": file_path}))
        except Exception:
            pass
    if not docs:
        try:
            docs = PyPDFLoader(file_path).load()
        except Exception:
            pass
    if not docs and OCR_AVAILABLE:
        try:
            images = convert_from_path(file_path)
            full_text = ""
            for i, img in enumerate(images):
                full_text += f"\n\nPage {i+1}:\n" + pytesseract.image_to_string(img, lang="eng+ara")
            if full_text.strip():
                docs = [Document(page_content=full_text, metadata={"source": file_path})]
        except Exception:
            pass
    return docs

def load_docx(file_path):
    try:
        return Docx2txtLoader(file_path).load()
    except Exception:
        return []

def load_pptx(file_path):
    if not PPTX_AVAILABLE:
        return []
    try:
        prs = PptxPresentation(file_path)
        pages = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text.strip() for cell in row.cells))
            if texts:
                pages.append(Document(page_content="\n".join(texts),
                                      metadata={"slide": i+1, "source": file_path}))
        return pages
    except Exception:
        return []

def load_excel(file_path):
    if not PANDAS_AVAILABLE:
        return []
    try:
        xl = pd.ExcelFile(file_path)
        docs = []
        for sheet in xl.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet).fillna("")
            docs.append(Document(
                page_content=f"Sheet: {sheet}\n" + df.to_string(index=False),
                metadata={"sheet": sheet, "source": file_path}
            ))
        return docs
    except Exception:
        return []

def load_csv(file_path):
    if not PANDAS_AVAILABLE:
        return []
    try:
        df = pd.read_csv(file_path).fillna("")
        return [Document(page_content=df.to_string(index=False),
                         metadata={"source": file_path})]
    except Exception:
        try:
            return CSVLoader(file_path).load()
        except Exception:
            return []

def load_image(file_path):
    if not OCR_AVAILABLE:
        return []
    try:
        text = pytesseract.image_to_string(Image.open(file_path), lang="eng+ara")
        if text.strip():
            return [Document(page_content=text, metadata={"source": file_path})]
    except Exception:
        pass
    return []

def load_file(file_path):
    ext = Path(file_path).suffix.lower()
    dispatch = {
        ".pdf":  load_pdf,  ".docx": load_docx, ".doc":  load_docx,
        ".pptx": load_pptx, ".ppt":  load_pptx,
        ".xlsx": load_excel,".xls":  load_excel, ".csv":  load_csv,
        ".txt":  lambda p: TextLoader(p, encoding="utf-8").load(),
        ".md":   lambda p: UnstructuredMarkdownLoader(p).load(),
        ".png":  load_image,".jpg":  load_image, ".jpeg": load_image, ".webp": load_image,
    }
    fn = dispatch.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return fn(file_path)

# ==============================================================================
# 🧹  TEXT CLEANING
# ==============================================================================

WATERMARKS = ["camscanner", "scanned by", "created with", "adobe scan", "watermark"]

def clean_text(text):
    low = text.lower()
    for wm in WATERMARKS:
        if wm in low and len(text) < 80:
            return ""
    text = re.sub(r'\s+', ' ', text).strip()
    return text if len(text) >= 25 and text.count(' ') >= 3 else ""

# ==============================================================================
# 🗄️  INDEXING
# ==============================================================================

def build_index(docs):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=900, chunk_overlap=180,
        separators=["\n\n", "\n", ". ", "! ", "? ", "، ", "؛ ", " ", ""],
    )
    chunks = splitter.split_documents(docs)
    chunks = [c for c in chunks if len(clean_text(c.page_content)) > 30]
    if not chunks:
        return None, None, []
    vector_db = FAISS.from_documents(chunks, get_embeddings())
    chunk_texts = [c.page_content for c in chunks]
    bm25 = None
    if BM25_AVAILABLE:
        bm25 = BM25Okapi([t.lower().split() for t in chunk_texts])
    return vector_db, bm25, chunk_texts

# ==============================================================================
# 🔍  HYBRID RETRIEVAL
# ==============================================================================

def hybrid_retrieve(query, doc_ids, top_k=6, doc_store=None):
    if not doc_ids:
        return []
    candidate_map = {}
    for doc_id in doc_ids:
        entry = (doc_store or {}).get(doc_id)
        if not entry:
            continue
        vdb         = entry["vector_db"]
        bm25        = entry.get("bm25")
        chunk_texts = entry["chunks"]
        source_name = entry["meta"]["name"]
        try:
            vec_results = vdb.similarity_search_with_score(query, k=top_k)
        except Exception:
            vec_results = []
        vec_rank = {}
        for rank, (doc, score) in enumerate(vec_results):
            key = doc.page_content[:120]
            vec_rank[key] = rank + 1
            candidate_map[key] = {
                "text": doc.page_content,
                "source": source_name,
                "rrf_score": 0.0,
                "vec_score": float(score),
                "page": doc.metadata.get("page", doc.metadata.get("slide", "?")),
            }
        bm25_rank = {}
        if bm25 and chunk_texts:
            try:
                scores = bm25.get_scores(query.lower().split())
                for rank, idx in enumerate(
                    sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:top_k]
                ):
                    key = chunk_texts[idx][:120]
                    bm25_rank[key] = rank + 1
                    if key not in candidate_map:
                        candidate_map[key] = {
                            "text": chunk_texts[idx], "source": source_name,
                            "rrf_score": 0.0, "vec_score": 0.0, "page": "?",
                        }
            except Exception:
                pass
        for key in set(vec_rank) | set(bm25_rank):
            rrf = 0.0
            if key in vec_rank:  rrf += 1.0 / (50 + vec_rank[key])
            if key in bm25_rank: rrf += 1.0 / (50 + bm25_rank[key])
            if key in candidate_map:
                candidate_map[key]["rrf_score"] += rrf

    candidates = sorted(candidate_map.values(),
                        key=lambda x: x["rrf_score"], reverse=True)[:top_k * 3]
    if not candidates:
        return []

    reranker = get_reranker()
    if reranker and len(candidates) > 1:
        try:
            re_scores = reranker.predict([(query, c["text"]) for c in candidates])
            for i, c in enumerate(candidates):
                c["rerank_score"] = float(re_scores[i])
            candidates.sort(key=lambda x: x.get("rerank_score", 0), reverse=True)
        except Exception:
            pass

    return candidates[:top_k]

# ==============================================================================
# 📄  DOCUMENT PROCESSING
# ==============================================================================

def file_hash(path: str) -> str:
    h = hashlib.md5()
    with open(path, "rb") as f:
        while True:
            chunk = f.read(8192)
            if not chunk:
                break
            h.update(chunk)
    return h.hexdigest()

def store_user_upload(file_path: str, username: str, session_id: Optional[str] = None):
    user_root = ensure_user_dirs(username, session_id=session_id)
    uploads_dir = user_root / "uploads" / (session_id or FALLBACK_SESSION_DIR)
    uploads_dir.mkdir(parents=True, exist_ok=True)
    target = uploads_dir / Path(file_path).name
    try:
        shutil.copy2(file_path, target)
    except Exception:
        pass

def process_files(files, username: str, session_id: Optional[str] = None):
    if not is_valid_username(username):
        return "⚠️ Please enter your name first."
    if not files:
        return "📁 Ready — upload one or more files to get started."

    sid = get_or_create_session_id(username, session_id)
    ensure_user_dirs(username, session_id=sid)
    clear_user_context(username, sid)
    ctx         = get_user_context(username, sid)
    doc_store   = ctx["doc_store"]
    active_docs = ctx["active_docs"]
    doc_hashes  = ctx["doc_hashes"]

    file_list = files if isinstance(files, list) else [files]
    results   = []
    loaded_names = []
    for file in file_list:
        file_path = file if isinstance(file, str) else file.name
        file_name = Path(file_path).name

        try:
            store_user_upload(file_path, username, sid)
            fhash = file_hash(file_path)
            if fhash in doc_hashes:
                results.append(f"⚠️ {file_name} — already indexed (duplicate).")
                continue
            doc_hashes.add(fhash)

            doc_id   = hashlib.md5(file_name.encode()).hexdigest()[:8]
            raw_docs = load_file(file_path)
            if not raw_docs:
                results.append(f"❌ {file_name} — could not extract text.")
                continue
            for d in raw_docs:
                d.page_content = clean_text(d.page_content)
            raw_docs = [d for d in raw_docs if d.page_content]
            if not raw_docs:
                results.append(f"❌ {file_name} — no usable text after cleaning.")
                continue

            vector_db, bm25, chunk_texts = build_index(raw_docs)
            if vector_db is None:
                results.append(f"❌ {file_name} — insufficient content.")
                continue

            doc_store[doc_id] = {
                "vector_db": vector_db, "bm25": bm25, "chunks": chunk_texts,
                "meta": {
                    "name":   file_name,
                    "pages":  len(raw_docs),
                    "chunks": len(chunk_texts),
                    "doc_id": doc_id,
                },
            }
            if doc_id not in active_docs:
                active_docs.append(doc_id)
            loaded_names.append(file_name)
            results.append(
                f"✅ {file_name}\n"
                f"   📊 {len(raw_docs)} pages → {len(chunk_texts)} chunks\n"
                f"   🔍 Hybrid {'+ BM25 ' if bm25 else ''}{'+ Reranker' if RERANKER_AVAILABLE else ''} ready"
            )
        except ValueError as e:
            results.append(f"⚠️ {file_name} — {e}")
        except Exception as e:
            results.append(f"❌ {file_name} — {e}")

    if loaded_names:
        update_session_docs(sid, [doc_store[d]["meta"]["name"]
                                   for d in active_docs if d in doc_store])

    loaded = [doc_store[d]["meta"]["name"] for d in active_docs if d in doc_store]
    return "\n\n".join(results) + f"\n\n---\n📚 Active ({len(loaded)}): {', '.join(loaded) or 'none'}"


def clear_documents(username: str, session_id: Optional[str] = None) -> str:
    if not is_valid_username(username):
        return "⚠️ Enter your name first."
    clear_user_context(username, session_id)
    return "🗑️ All documents cleared. Upload new files to start fresh."

def list_documents(username: str = "", session_id: Optional[str] = None):
    ctx = get_user_context(username, session_id) if username else {"doc_store": {}, "active_docs": []}
    doc_store   = ctx["doc_store"]
    active_docs = ctx["active_docs"]
    if not doc_store:
        return "No documents indexed yet."
    lines = []
    for doc_id, entry in doc_store.items():
        m = entry["meta"]
        lines.append(f"{'🟢' if doc_id in active_docs else '🔴'} {m['name']} — {m['chunks']} chunks")
    return "\n".join(lines)

# ==============================================================================
# 🌐  LANGUAGE DETECTION + DOC REQUIRED
# ==============================================================================

DOC_REQUEST_PATTERNS = re.compile(
    r"(summarize|summary|tldr|tl;dr|explain.*document|analyze|analysis|"
    r"compare|mind map|flashcard|quiz|extract|notes|outline|"
    r"document|pdf|file|slides|ppt|chapter|sections|pages|"
    r"الملف|الوثيقة|المستند|الملفات|المستندات|"
    r"الشرائح|المحاضرة|الكتاب|الفصل|الصفحة|"
    r"لخص|ملخص|تحليل|قارن|استخرج|بيانات)",
    re.I
)

def detect_language(text):
    arabic = len(re.findall(r'[\u0600-\u06FF]', text))
    total  = len(re.findall(r'[a-zA-Z\u0600-\u06FF]', text))
    return "ar" if total and (arabic / total) > 0.3 else "en"

def requires_document(message: str, mode: str) -> bool:
    if mode in {"quiz", "flashcard", "summary", "compare", "mindmap"}:
        return True
    return bool(DOC_REQUEST_PATTERNS.search(message))

# ==============================================================================
# 🤖  QUERY REWRITING + SUMMARIES
# ==============================================================================

FOLLOWUP_PATTERNS = re.compile(
    r"^(what about|how about|and|also|"
    r"tell me more|elaborate|expand|explain that|"
    r"why|which one|can you|could you|"
    r"ما عن|اشرح|وضح|أكمل|ولكن|ماذا عن)",
    re.I
)

def is_followup(message: str, history: list) -> bool:
    if not history:
        return False
    if len(message.split()) <= 5 and FOLLOWUP_PATTERNS.match(message.strip()):
        return True
    return False

def build_contextual_query(message: str, history: list) -> str:
    if not history or not is_followup(message, history):
        return message
    last_assistant = ""
    for turn in reversed(history):
        if isinstance(turn, dict) and turn.get("role") == "assistant":
            last_assistant = turn.get("content", "")[:300]
            break
    if last_assistant:
        return f"{last_assistant}\n\nFollow-up question: {message}"
    return message

def rewrite_query(question: str, lang: str) -> str:
    if lang == "ar":
        prompt = f"أعد صياغة هذا السؤال كاستعلام بحث دقيق. اكتب الاستعلام فقط:\n{question}"
    else:
        prompt = (
            f"Rewrite the following as a precise document search query. "
            f"Output only the query, no preamble:\n{question}"
        )
    try:
        resp = client.chat_completion(
            [{"role": "user", "content": prompt}],
            max_tokens=80, temperature=0.2,
        )
        rewritten = resp.choices[0].message.content.strip()
        return rewritten if len(rewritten) > 5 else question
    except Exception:
        return question

def update_summary_if_needed(session_id: str):
    assistant_turns = count_assistant_turns(session_id)
    if assistant_turns == 0 or assistant_turns % SUMMARY_EVERY_N_ASSISTANT_TURNS != 0:
        return
    summary = get_session_summary(session_id)
    recent = load_session_messages(session_id, limit=12)
    history_text = "\n".join(
        [f"{'User' if m['role']=='user' else 'Assistant'}: {m['content'][:300]}" for m in recent]
    )
    prompt = (
        "Update the ongoing session summary. Keep it short, factual, and useful.\n"
        "If no prior summary, create one.\n\n"
        f"PREVIOUS SUMMARY:\n{summary}\n\n"
        f"RECENT DIALOGUE:\n{history_text}\n\n"
        "NEW SUMMARY (max 120 words):"
    )
    try:
        resp = client.chat_completion(
            [{"role": "user", "content": prompt}],
            max_tokens=200, temperature=0.2,
        )
        new_summary = resp.choices[0].message.content.strip()
        if new_summary:
            update_session_summary(session_id, new_summary)
    except Exception:
        pass

# ==============================================================================
# 🧠  MODE DETECTION + SYSTEM PROMPTS
# ==============================================================================

QUIZ_TRIGGER    = re.compile(r'(quiz|mcq|test me|generate questions|أسئلة اختيار|اختبار)', re.I)
FLASH_TRIGGER   = re.compile(r'(flashcard|flash card|بطاقة|بطاقات دراسية|flash)', re.I)
SUMMARY_TRIGGER = re.compile(r'(summarize|summary|tldr|tl;dr|ملخص|لخص|summarise)', re.I)
COMPARE_TRIGGER = re.compile(r'(compare|versus|\bvs\b|difference between|قارن|مقارنة)', re.I)
MINDMAP_TRIGGER = re.compile(r'(mind ?map|concept map|خريطة ذهنية)', re.I)
EXPLAIN_TRIGGER = re.compile(r'(explain.*simple|explain.*kid|eli5|explain.*15|اشرح.*بسيط)', re.I)

def detect_mode(message: str) -> str:
    if QUIZ_TRIGGER.search(message):    return "quiz"
    if FLASH_TRIGGER.search(message):   return "flashcard"
    if SUMMARY_TRIGGER.search(message): return "summary"
    if COMPARE_TRIGGER.search(message): return "compare"
    if MINDMAP_TRIGGER.search(message): return "mindmap"
    if EXPLAIN_TRIGGER.search(message): return "simple"
    return "chat"

def build_system_prompt(lang: str, mode: str, has_docs: bool,
                        history_context: str = "", summary: str = "") -> str:
    doc_instruction = (
        "You have access to retrieved document excerpts below. "
        "When answering, always cite the source name and page number using the format "
        "[Source: filename, Page X]. If multiple chunks support the answer, cite all of them. "
        "If the answer is not in the documents, say so clearly and answer from general knowledge."
        if has_docs else
        "No documents are indexed. Answer from your general knowledge and be transparent about it."
    )
    history_section = (
        f"\n\n### Conversation summary (context only):\n{summary}\n"
        if summary else ""
    ) + (
        f"\n\n### Recent context (do NOT re-answer these):\n{history_context}\n"
        if history_context else ""
    )
    if lang == "ar":
        base = (
            f"أنت Phoenix، مساعد دراسي ذكي ومتقدم.\n"
            f"{doc_instruction}\n"
            f"{history_section}"
            "القواعد:\n"
            "1. استشهد بالمصدر والصفحة عند كل معلومة مستخرجة من المستند.\n"
            "2. كن دقيقاً ومنظماً ومفيداً.\n"
            "3. إذا لم تجد الإجابة في المستند، قل ذلك صراحةً ثم أجب من معرفتك.\n"
            "4. لا تختلق اقتباسات أو مصادر.\n"
        )
        extras = {
            "quiz": (
                "\nأنشئ 5 أسئلة اختيار من متعدد من السياق. "
                "استخدم هذا التنسيق بالضبط لكل سؤال:\n"
                "**س{N}.** [السؤال]\n"
                "أ) ...\n ب) ...\n ج) ...\n د) ...\n"
                "✅ **الإجابة:** [الحرف] — [تفسير مختصر]\n\n"
            ),
            "flashcard": (
                "\nأنشئ 6 بطاقات دراسية. استخدم هذا التنسيق لكل بطاقة:\n"
                "---\n🃏 **الأمامي:** [المفهوم أو السؤال]\n"
                "📖 **الخلفي:** [الشرح أو الإجابة]\n"
            ),
            "summary": (
                "\nقدم ملخصاً منظماً:\n"
                "## نظرة عامة\n[فقرة مختصرة]\n\n"
                "## النقاط الرئيسية\n- ...\n\n"
                "## المصطلحات المهمة\n- **المصطلح**: التعريف\n"
            ),
            "compare": "\nقدم مقارنة واضحة بتنسيق جدول أو أقسام مع أوجه التشابه والاختلاف.\n",
            "mindmap": (
                "\nأنشئ خريطة ذهنية نصية هرمية باستخدام المسافة البادئة:\n"
                "# [الموضوع الرئيسي]\n"
                "## [فرع 1]\n   - [تفاصيل]\n"
                "## [فرع 2]\n   - [تفاصيل]\n"
            ),
            "simple": "\nاشرح هذا الموضوع بأسلوب بسيط جداً كأنك تشرح لطالب في الخامسة عشرة.\n",
        }
    else:
        base = (
            f"You are Phoenix, an advanced AI study assistant.\n"
            f"{doc_instruction}\n"
            f"{history_section}"
            "Rules:\n"
            "1. Cite sources inline as [Source: filename, Page X].\n"
            "2. If the question is a follow-up, refer back to the earlier answer.\n"
            "3. Be precise, structured, and genuinely helpful.\n"
            "4. Never fabricate citations.\n"
        )
        extras = {
            "quiz": (
                "\nGenerate exactly 5 multiple-choice questions from the context. "
                "Format each question exactly as:\n"
                "**Q{N}.** [question]\n"
                "A) ...\nB) ...\nC) ...\nD) ...\n"
                "✅ **Answer:** [letter] — [brief explanation]\n\n"
            ),
            "flashcard": (
                "\nGenerate exactly 6 flashcards from the content. "
                "Format each card exactly as:\n\n"
                "---\n"
                "🃏 **Front:** [concept or question]\n"
                "📖 **Back:** [clear explanation or answer]\n\n"
                "Repeat for all 6 cards."
            ),
            "summary": (
                "\nProvide a structured summary:\n"
                "## Overview\n[2-3 sentence overview]\n\n"
                "## Key Points\n- ...\n\n"
                "## Important Terms\n- **term**: definition\n"
            ),
            "compare": (
                "\nCreate a clear comparison with Similarities and Differences sections. "
                "Use a markdown table where possible.\n"
            ),
            "mindmap": (
                "\nCreate a hierarchical text mind map using indentation:\n"
                "# [Main Topic]\n"
                "## [Branch 1]\n   - detail\n"
                "## [Branch 2]\n   - detail\n"
            ),
            "simple": (
                "\nExplain this as if talking to a curious 15-year-old. "
                "Use analogies, examples, and plain language. Avoid jargon.\n"
            ),
        }
    return base + extras.get(mode, "")

# ==============================================================================
# 💬  CHAT LOGIC
# ==============================================================================

def chat_logic(message: str, history: list, username: str, session_id: Optional[str] = None) -> Generator[str, None, None]:
    if not is_valid_username(username):
        yield "⚠️ Please enter your name first."
        return
    if not message or not message.strip():
        yield "Please ask a question. | يرجى طرح سؤال."
        return

    lang    = detect_language(message)
    mode    = detect_mode(message)
    sid     = get_or_create_session_id(username, session_id)
    uname   = sanitize_username(username)

    save_message(sid, "user", message.strip())
    append_history_file(uname, sid, "user", message.strip())

    ctx         = get_user_context(username, sid)
    doc_store   = ctx["doc_store"]
    active_docs = ctx["active_docs"]

    if not active_docs and requires_document(message, mode):
        reply = "Hi, I'm Phoenix. Please upload a document first so I can help you with that."
        save_message(sid, "assistant", reply)
        append_history_file(uname, sid, "assistant", reply)
        yield reply
        return

    db_history = load_session_messages(sid, limit=MAX_HISTORY_TURNS)
    history_lines = []
    for m in db_history[-10:]:
        prefix = "User" if m["role"] == "user" else "Phoenix"
        history_lines.append(f"{prefix}: {m['content'][:350]}")
    history_context = "\n".join(history_lines[:-1])

    source_refs = []
    has_docs    = bool(active_docs)
    if has_docs:
        contextual_q = build_contextual_query(message, history)
        search_query = rewrite_query(contextual_q, lang) if len(message.split()) > 3 else message
        results      = hybrid_retrieve(search_query, active_docs, top_k=6, doc_store=doc_store)
        if results:
            sections = []
            for i, r in enumerate(results):
                page = r.get("page", "?")
                src  = r.get("source", "doc")
                label = f"[Source: {src} | Page {page} | Chunk {i+1}]"
                sections.append(f"{label}:\n{r['text']}")
                source_refs.append({
                    "source": src, "page": str(page),
                    "snippet": r["text"][:120],
                })
            context_block = "\n\n".join(sections)
            context_block = context_block[:MAX_CONTEXT_CHARS]
        else:
            context_block = ""
            has_docs = False
    else:
        context_block = ""

    summary    = get_session_summary(sid)
    system_msg = build_system_prompt(lang, mode, has_docs, history_context, summary)
    if context_block:
        system_msg += f"\n\n--- DOCUMENT CONTEXT ---\n{context_block}\n--- END CONTEXT ---"

    messages = [{"role": "system", "content": system_msg}]
    recent = history[-6:] if len(history) > 6 else history
    for turn in recent:
        role    = turn.get("role")    if isinstance(turn, dict) else None
        content = turn.get("content") if isinstance(turn, dict) else None
        if role in ("user", "assistant") and content:
            messages.append({"role": role, "content": str(content)})
    messages.append({"role": "user", "content": message.strip()})

    partial     = ""
    buffer      = ""
    final_reply = ""
    try:
        stream = client.chat_completion(
            messages,
            max_tokens=1800,
            stream=True,
            temperature=0.35,
            top_p=0.9,
        )
        for chunk in stream:
            if chunk.choices and chunk.choices[0].delta.content:
                raw = chunk.choices[0].delta.content
                if isinstance(raw, list):
                    raw = "".join(item.get("text","") if isinstance(item,dict) else str(item) for item in raw)
                elif not isinstance(raw, str):
                    raw = str(raw)
                buffer += raw
                while " " in buffer:
                    word, buffer = buffer.split(" ", 1)
                    partial += word + " "
                    yield partial
        if buffer:
            partial += buffer
            yield partial
        final_reply = partial
    except Exception as e:
        err = str(e)
        if "410" in err or "deprecated" in err.lower():
            final_reply = "🔴 Model unavailable. Check REPO_ID."
        elif "503" in err:
            final_reply = "🟡 Server loading — please wait 60 s and retry."
        elif "401" in err or "403" in err:
            final_reply = "🔑 Auth error — check HF_TOKEN."
        else:
            final_reply = f"⚠️ {err}"
        yield final_reply

    if final_reply:
        save_message(sid, "assistant", final_reply, sources=source_refs, confidence=0.0)
        append_history_file(uname, sid, "assistant", final_reply)
        update_summary_if_needed(sid)

# ==============================================================================
# 🎧  TEXT → AUDIO + VIDEO
# ==============================================================================

VOICE_MAP = {
    "English": "en-US-JennyNeural",
    "Arabic":  "ar-EG-SalmaNeural",
}

def resolve_lang(choice: str, text: str) -> str:
    if choice == "Auto":
        return detect_language(text or "test")
    return "ar" if choice == "Arabic" else "en"

def tts_to_audio(text: str, lang: str, out_path: str):
    if AZURE_TTS_AVAILABLE and AZURE_SPEECH_KEY and AZURE_SPEECH_REGION:
        voice = VOICE_MAP["Arabic"] if lang == "ar" else VOICE_MAP["English"]
        speech_config = speechsdk.SpeechConfig(
            subscription=AZURE_SPEECH_KEY, region=AZURE_SPEECH_REGION
        )
        speech_config.speech_synthesis_voice_name = voice
        audio_config  = speechsdk.audio.AudioOutputConfig(filename=out_path)
        synthesizer   = speechsdk.SpeechSynthesizer(
            speech_config=speech_config, audio_config=audio_config
        )
        synthesizer.speak_text_async(text).get()
        return out_path

    if GTTS_AVAILABLE:
        tts = gTTS(text=text, lang="ar" if lang == "ar" else "en")
        tts.save(out_path)
        return out_path

    raise ValueError("No TTS engine available. Install azure-cognitiveservices-speech or gTTS.")

"""
DROP-IN REPLACEMENT for the video/slide section of app.py
Replace everything from  chunk_text()  through  text_to_video()
with this block.
"""

import math
import textwrap

# ── Font paths (guaranteed available on HF Spaces) ───────────────────────────
_FONT_BOLD    = "/usr/share/fonts/truetype/google-fonts/Poppins-Bold.ttf"
_FONT_MEDIUM  = "/usr/share/fonts/truetype/google-fonts/Poppins-Medium.ttf"
_FONT_REGULAR = "/usr/share/fonts/truetype/google-fonts/Poppins-Regular.ttf"
_FONT_LIGHT   = "/usr/share/fonts/truetype/google-fonts/Poppins-Light.ttf"
_FONT_FALLBACK = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"


def _font(path: str, size: int):
    from PIL import ImageFont
    try:
        return ImageFont.truetype(path, size)
    except Exception:
        try:
            return ImageFont.truetype(_FONT_FALLBACK, size)
        except Exception:
            return ImageFont.load_default()


def chunk_text(text: str, max_chars: int = 380):
    """Split narration text into slide-sized chunks, breaking on sentences."""
    import re
    # Split on sentence boundaries first
    sentences = re.split(r'(?<=[.!?؟])\s+', text.strip())
    chunks, buf = [], ""
    for s in sentences:
        if len(buf) + len(s) + 1 <= max_chars:
            buf += s + " "
        else:
            if buf.strip():
                chunks.append(buf.strip())
            # Long single sentence → hard wrap
            if len(s) > max_chars:
                for part in textwrap.wrap(s, max_chars):
                    chunks.append(part)
                buf = ""
            else:
                buf = s + " "
    if buf.strip():
        chunks.append(buf.strip())
    return chunks or [text[:max_chars]]


# ── Colour palette (matches Phoenix dark theme) ───────────────────────────────
_BG_TOP    = (10,  11,  18)
_BG_BOT    = (18,  16,  28)
_ACCENT    = (255, 107,  61)   # Phoenix orange
_ACCENT2   = (255, 160, 100)
_GOLD      = (240, 195,  70)
_TEXT_HEAD = (255, 230, 210)
_TEXT_BODY = (220, 220, 235)
_TEXT_DIM  = (130, 140, 165)
_CARD_BG   = (22,  22,  35)
_DIVIDER   = (50,  42,  68)


def _draw_gradient_bg(img):
    """Vertical gradient background."""
    from PIL import Image
    w, h = img.size
    pix  = img.load()
    for y in range(h):
        t = y / h
        r = int(_BG_TOP[0] + t * (_BG_BOT[0] - _BG_TOP[0]))
        g = int(_BG_TOP[1] + t * (_BG_BOT[1] - _BG_TOP[1]))
        b = int(_BG_TOP[2] + t * (_BG_BOT[2] - _BG_TOP[2]))
        for x in range(w):
            pix[x, y] = (r, g, b)


def _draw_rounded_rect(draw, xy, radius, fill, outline=None, outline_width=2):
    """Draw a rounded rectangle (Pillow ≥ 9.2 has this built-in; fallback for older)."""
    x0, y0, x1, y1 = xy
    try:
        draw.rounded_rectangle(xy, radius=radius, fill=fill,
                                outline=outline, width=outline_width)
    except AttributeError:
        draw.rectangle(xy, fill=fill, outline=outline, width=outline_width)


def _wrap_text_to_lines(text: str, font, max_width: int, draw) -> list:
    """Word-wrap text to fit max_width pixels."""
    words  = text.split()
    lines  = []
    line   = ""
    for word in words:
        test = (line + " " + word).strip()
        try:
            w = draw.textlength(test, font=font)
        except AttributeError:
            w = font.getlength(test)
        if w <= max_width:
            line = test
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def create_slide_image(text: str, out_path: str,
                       title: str = "Phoenix Eduplan",
                       slide_num: int = 1, total_slides: int = 1):
    """
    Render a polished Phoenix-branded slide image.
    No OCR required — only Pillow.
    """
    from PIL import Image, ImageDraw

    W, H = 1280, 720
    img  = Image.new("RGB", (W, H))
    _draw_gradient_bg(img)
    draw = ImageDraw.Draw(img)

    # ── Decorative top bar ────────────────────────────────────────────────────
    draw.rectangle([(0, 0), (W, 5)], fill=_ACCENT)

    # ── Subtle grid pattern overlay ───────────────────────────────────────────
    for x in range(0, W, 80):
        draw.line([(x, 0), (x, H)], fill=(255, 255, 255, 8), width=1)
    for y in range(0, H, 80):
        draw.line([(0, y), (W, y)], fill=(255, 255, 255, 8), width=1)

    # ── Logo / brand strip (top-left) ─────────────────────────────────────────
    logo_font  = _font(_FONT_BOLD, 22)
    draw.text((36, 22), "🦅 PHOENIX EDUPLAN", font=logo_font, fill=_ACCENT)

    # ── Slide counter (top-right) ─────────────────────────────────────────────
    ctr_font = _font(_FONT_REGULAR, 18)
    ctr_txt  = f"{slide_num} / {total_slides}"
    try:
        ctr_w = draw.textlength(ctr_txt, font=ctr_font)
    except AttributeError:
        ctr_w = ctr_font.getlength(ctr_txt)
    draw.text((W - ctr_w - 36, 26), ctr_txt, font=ctr_font, fill=_TEXT_DIM)

    # ── Thin divider line ─────────────────────────────────────────────────────
    draw.rectangle([(36, 58), (W - 36, 60)], fill=_DIVIDER)

    # ── Content card ─────────────────────────────────────────────────────────
    card_x0, card_y0 = 54, 76
    card_x1, card_y1 = W - 54, H - 70
    _draw_rounded_rect(draw,
                       (card_x0, card_y0, card_x1, card_y1),
                       radius=18,
                       fill=_CARD_BG,
                       outline=_DIVIDER, outline_width=1)

    # ── Accent dot cluster (decorative, top-right of card) ────────────────────
    for i, col in enumerate([_ACCENT, _GOLD, (100, 200, 255)]):
        cx = card_x1 - 28 - i * 20
        cy = card_y0 + 22
        draw.ellipse([(cx - 6, cy - 6), (cx + 6, cy + 6)], fill=col)

    # ── Decide layout: short vs long text ────────────────────────────────────
    padding   = 52
    text_x    = card_x0 + padding
    max_w     = (card_x1 - card_x0) - padding * 2
    text_area_top = card_y0 + padding

    word_count = len(text.split())

    if word_count <= 18:
        # ── BIG centred quote style ───────────────────────────────────────────
        # Orange accent line left
        draw.rectangle([(card_x0 + 24, card_y0 + 90),
                        (card_x0 + 30, card_y1 - 90)],
                       fill=_ACCENT)

        body_font  = _font(_FONT_BOLD, 42)
        lines      = _wrap_text_to_lines(text, body_font, max_w - 20, draw)
        line_h     = 56
        total_h    = len(lines) * line_h
        y          = (H - total_h) // 2 - 10
        for line in lines:
            draw.text((text_x + 16, y), line, font=body_font, fill=_TEXT_HEAD)
            y += line_h

    else:
        # ── Multi-line body text ──────────────────────────────────────────────
        body_font  = _font(_FONT_REGULAR, 30)
        lines      = _wrap_text_to_lines(text, body_font, max_w, draw)
        line_h     = 44
        max_lines  = int((card_y1 - card_y0 - padding * 2) / line_h)
        lines      = lines[:max_lines]

        # Leading bullet dots
        bullet_x   = text_x - 20
        y          = text_area_top + 10
        for line in lines:
            # small accent dot per paragraph-start heuristic
            draw.ellipse([(bullet_x - 4, y + 14),
                          (bullet_x + 4, y + 22)],
                         fill=_ACCENT)
            draw.text((text_x, y), line, font=body_font, fill=_TEXT_BODY)
            y += line_h

    # ── Bottom bar ────────────────────────────────────────────────────────────
    draw.rectangle([(0, H - 52), (W, H)], fill=(14, 12, 22))
    bar_font = _font(_FONT_LIGHT, 17)
    draw.text((36, H - 34), "Study with Phoenix · AI-powered learning",
              font=bar_font, fill=_TEXT_DIM)
    # Progress bar
    prog_w = int((W - 72) * slide_num / total_slides)
    draw.rectangle([(36, H - 6), (36 + prog_w, H - 2)], fill=_ACCENT)
    draw.rectangle([(36 + prog_w, H - 6), (W - 36, H - 2)], fill=_DIVIDER)

    img.save(out_path, quality=95)
    return out_path


def text_to_video(text: str, lang: str) -> str:
    import os, shutil, tempfile
    if not MOVIEPY_AVAILABLE:
        raise ValueError(
            "moviepy is not installed. Add 'moviepy==1.0.3' and 'decorator<5.0' "
            "to requirements.txt and restart the Space."
        )

    with tempfile.TemporaryDirectory() as tmp:
        # 1. Generate full narration audio first
        audio_path = os.path.join(tmp, "narration.mp3")
        tts_to_audio(text, lang, audio_path)

        # 2. Split into slide-sized chunks
        parts = chunk_text(text, max_chars=380)
        total = len(parts)

        # 3. Render each slide
        clips = []
        for i, part in enumerate(parts):
            img_path = os.path.join(tmp, f"slide_{i:03d}.png")
            create_slide_image(part, img_path,
                               slide_num=i + 1,
                               total_slides=total)
            # Estimate duration proportional to word count
            words    = len(part.split())
            duration = max(3.5, words * 0.45)   # ~0.45 s per word, min 3.5 s
            clips.append(ImageClip(img_path).set_duration(duration))

        # 4. Assemble video with narration
        from moviepy.editor import AudioFileClip, concatenate_videoclips
        audio     = AudioFileClip(audio_path)
        video_raw = concatenate_videoclips(clips, method="compose")

        # Trim/loop to match audio length exactly
        if audio.duration > video_raw.duration:
            # Extend last slide to cover remaining audio
            extra = audio.duration - video_raw.duration + 0.1
            last  = clips[-1].set_duration(clips[-1].duration + extra)
            clips[-1] = last
            video_raw = concatenate_videoclips(clips, method="compose")

        video = video_raw.set_audio(
            audio.subclip(0, min(audio.duration, video_raw.duration))
        )

        out_path = os.path.join(tmp, "study_video.mp4")
        video.write_videofile(
            out_path,
            fps=24,
            codec="libx264",
            audio_codec="aac",
            bitrate="2500k",          # higher = sharper slides
            audio_bitrate="192k",
            preset="fast",
            verbose=False,
            logger=None,
        )

        final_path = os.path.join(
            os.getcwd(),
            f"study_video_{int(datetime.utcnow().timestamp())}.mp4"
        )
        shutil.copy(out_path, final_path)
        return final_path

# FIX 2: Strict script-generation prompt — no preamble, no meta-commentary
def build_script_from_docs(username: str, user_text: str, lang: str, session_id: Optional[str] = None) -> str:
    ctx         = get_user_context(username, session_id)
    doc_store   = ctx["doc_store"]
    active_docs = ctx["active_docs"]

    if not active_docs:
        return user_text.strip()

    query = user_text.strip() if user_text.strip() else (
        "Summarize the documents into a short narration script." if lang == "en"
        else "لخّص المستندات كنص سردي قصير."
    )

    results = hybrid_retrieve(query, active_docs, top_k=6, doc_store=doc_store)
    if not results:
        return user_text.strip()

    context = "\n\n".join([r["text"] for r in results])[:3500]

    # ── STRICT: output ONLY the spoken narration, zero preamble ──
    if lang == "en":
        prompt = (
            "You are a narrator. Output ONLY the spoken narration script — "
            "no introduction, no 'here is a summary', no meta-commentary, "
            "no labels, no headings. Just the words to be spoken aloud.\n\n"
            f"User request: {user_text or 'Narrate the key ideas from the document.'}\n\n"
            f"Document context:\n{context}"
        )
    else:
        prompt = (
            "أنت راوٍ. اكتب نص السرد المنطوق فقط — "
            "بدون مقدمة، بدون 'إليك ملخصاً'، بدون تعليق، بدون عناوين. "
            "فقط الكلمات التي ستُقرأ بصوت عالٍ.\n\n"
            f"طلب المستخدم: {user_text or 'اسرد الأفكار الرئيسية من المستند.'}\n\n"
            f"سياق المستند:\n{context}"
        )

    try:
        resp = client.chat_completion(
            [{"role": "user", "content": prompt}],
            max_tokens=600,
            temperature=0.3,
        )
        result = resp.choices[0].message.content.strip()
        # Defensive: strip any leading meta-line that slipped through
        lines = result.splitlines()
        clean_lines = []
        for line in lines:
            low = line.lower().strip()
            if any(bad in low for bad in [
                "here is", "here's", "below is", "the following",
                "إليك", "فيما يلي", "هذا هو", "هذا ملخص",
            ]):
                continue
            clean_lines.append(line)
        return "\n".join(clean_lines).strip() or user_text.strip()
    except Exception:
        return user_text.strip()


def generate_audio_only(text: str, username: str, use_docs: bool, lang_choice: str, session_id: Optional[str] = None):
    try:
        if not text.strip() and not use_docs:
            return None, "⚠️ Enter text or enable 'Use uploaded docs'."
        lang   = resolve_lang(lang_choice, text)
        script = build_script_from_docs(username, text, lang, session_id) if use_docs else text
        audio_path = os.path.join(
            os.getcwd(), f"audio_{int(datetime.utcnow().timestamp())}.mp3"
        )
        tts_to_audio(script, lang, audio_path)
        return audio_path, "✅ Audio generated."
    except Exception as e:
        return None, f"⚠️ Audio error: {e}"

def generate_video_only(text: str, username: str, use_docs: bool, lang_choice: str, session_id: Optional[str] = None):
    try:
        if not text.strip() and not use_docs:
            return None, "⚠️ Enter text or enable 'Use uploaded docs'."
        lang   = resolve_lang(lang_choice, text)
        script = build_script_from_docs(username, text, lang, session_id) if use_docs else text
        video_path = text_to_video(script, lang)
        return video_path, "✅ Video generated."
    except Exception as e:
        return None, f"⚠️ Video error: {e}"

# ==============================================================================
# 🎨  CSS + UI TEXT
# ==============================================================================

DARK_CSS = """
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;700;800&family=DM+Sans:wght@300;400;500&display=swap');
:root {
    --bg0:#080a10;
    --bg1:#0f1117;
    --bg2:#161820;
    --bg3:#1d2030;
    --border:#252838;
    --border2:#353850;
    --accent:#ff6b3d;
    --accent2:#cc4b28;
    --accentg:rgba(255,107,61,0.18);
    --gold:#f0c040;
    --green:#22d47a;
    --red:#f87171;
    --t1:#eceef5;
    --t2:#8892aa;
    --t3:#4a5268;
    --r:12px;
}
body,html,.gradio-container,.gradio-container>div,
#root,[data-testid="block"],.gap,.contain,.wrap,div.prose{
    background-color:var(--bg0)!important;
    color:var(--t1)!important;
}
.gradio-container .block,.gradio-container fieldset,
.gradio-container [data-testid="column"],
.gradio-container .tabs,.gradio-container .tabitem{
    background:var(--bg1)!important;
    border-color:var(--border)!important;
}
.gradio-container textarea,.gradio-container input{
    background:var(--bg2)!important;
    color:var(--t1)!important;
    border:1px solid var(--border2)!important;
    border-radius:10px!important;
    font-family:'DM Sans',sans-serif!important;
    font-size:0.95rem!important;
    caret-color:var(--accent)!important;
}
.gradio-container textarea:focus,.gradio-container input:focus{
    border-color:var(--accent)!important;
    box-shadow:0 0 0 3px var(--accentg)!important;
    outline:none!important;
}
.gradio-container textarea::placeholder,
.gradio-container input::placeholder{
    color:var(--t3)!important;
}
.gradio-container .tab-nav,
.gradio-container [role="tablist"]{
    background:var(--bg1)!important;
    border-bottom:1px solid var(--border)!important;
    display:flex!important;
    flex-wrap:wrap!important;
    gap:2px!important;
}
.gradio-container [role="tab"]{
    flex:1 1 auto!important;
    white-space:nowrap!important;
    font-size:0.78rem!important;
    padding:6px 8px!important;
}
.gradio-container [data-testid="file-upload"]{
    background:var(--bg2)!important;
    border:2px dashed var(--border2)!important;
    border-radius:var(--r)!important;
    color:var(--t2)!important;
}
.gradio-container button{
    font-family:'Syne',sans-serif!important;
    font-weight:600!important;
    border-radius:10px!important;
    transition:all 0.2s!important;
    cursor:pointer!important;
}
.gradio-container button.primary,
.gradio-container button[variant="primary"]{
    background:var(--accent)!important;
    color:#fff!important;
    border:none!important;
}
.gradio-container .chatbot,
.gradio-container [data-testid="chatbot"]{
    background: linear-gradient(180deg, rgba(20,8,8,0.96) 0%, rgba(10,10,14,0.98) 100%)!important;
    border:1px solid rgba(255,110,60,0.18)!important;
    border-radius:var(--r)!important;
    position:relative;
    overflow:hidden;
}
footer,.footer{display:none!important;}

/* ── Password status badge ── */
.pw-status-box{
    margin-top:8px;
    padding:9px 14px;
    border-radius:10px;
    font-family:'Syne',sans-serif;
    font-size:0.88rem;
    font-weight:700;
    text-align:center;
    transition: all 0.3s ease;
}
.pw-locked{
    background:rgba(248,113,113,0.10);
    color:#fca5a5;
    border:1px solid rgba(248,113,113,0.25);
}
.pw-ok{
    background:rgba(34,212,122,0.12);
    color:#86efac;
    border:1px solid rgba(34,212,122,0.28);
}
/* ── Unlock button glow on success ── */
.unlock-btn-ok button{
    background:var(--green)!important;
    color:#000!important;
    box-shadow:0 0 18px rgba(34,212,122,0.45)!important;
}

.ph-header.flame{
    background: linear-gradient(180deg, rgba(20, 8, 8, 0.96) 0%, rgba(10, 10, 14, 0.98) 100%);
    border: 1px solid rgba(255, 110, 60, 0.18);
    box-shadow: 0 10px 30px rgba(0, 0, 0, 0.35);
    padding: 22px 28px;
    border-radius: 16px;
    margin: 14px 14px 18px 14px;
    position: relative;
    overflow: hidden;
}
.ph-title{
    font-family:'Syne',sans-serif;
    font-size:2.15rem;
    font-weight:800;
    letter-spacing:-0.03em;
    line-height:1.05;
    background: linear-gradient(90deg, #ffcf9f 0%, #ff7a3d 35%, #ff4d4d 70%, #ffd9b8 100%);
    -webkit-background-clip:text;
    -webkit-text-fill-color:transparent;
    background-clip:text;
}
.ph-version{
    font-size:0.52em;
    opacity:0.65;
    font-weight:600;
    margin-left:6px;
}
.ph-sub{
    color:#d8c7bf;
    font-size:0.95rem;
    margin-top:5px;
    font-weight:400;
    opacity:0.9;
}
.ph-pills{
    margin-top:14px;
    display:flex;
    gap:8px;
    flex-wrap:wrap;
}
.pill{
    display:inline-flex;
    align-items:center;
    gap:5px;
    padding:5px 12px;
    border-radius:999px;
    font-size:0.74rem;
    font-weight:700;
    font-family:'Syne',sans-serif;
    letter-spacing:0.01em;
}
.pill-red{
    background:rgba(255,80,80,0.10);
    color:#ffb2b2;
    border:1px solid rgba(255,90,90,0.22);
}
.pill-blue{
    background:rgba(108,114,255,0.12);
    color:#c5cbff;
    border:1px solid rgba(108,114,255,0.22);
}
.pill-green{
    background:rgba(34,212,122,0.10);
    color:#b5f5cf;
    border:1px solid rgba(34,212,122,0.22);
}
.pill-gold{
    background:rgba(240,192,64,0.10);
    color:#ffe6b0;
    border:1px solid rgba(240,192,64,0.20);
}
"""

NOTEBOOKLM_CSS = """
.gradio-container .tabitem {
    box-shadow: inset 0 0 0 1px rgba(255,255,255,0.04);
}
.gradio-container .block {
    border-radius: 14px !important;
}
.gradio-container .markdown-body, .gradio-container .prose {
    color: #d9dbe4 !important;
}
"""

HEADER_HTML = """
<div class="ph-header flame">
  <div class="ph-title">Phoenix Eduplan <span class="ph-version">v3.0</span></div>
  <div class="ph-sub">Advanced AI study assistant · Persistent storage · Multi-doc RAG · Hybrid retrieval</div>
  <div class="ph-pills">
    <span class="pill pill-red">🦅 Phoenix Core</span>
    <span class="pill pill-blue">⚡ Llama 3.3 70B</span>
    <span class="pill pill-green">🔍 Hybrid Search</span>
    <span class="pill pill-gold">📄 PDFs · DOCX · PPTX · XLSX · CSV · Images</span>
  </div>
</div>
"""

SIDEBAR_GUIDE = """
### 💡 Quick Commands
| Say... | Gets you... |
|--------|------------|
| `summarize` | Structured summary |
| `quiz me` | 5 MCQ questions |
| `flashcards` | 6 study cards |
| `compare X and Y` | Side-by-side comparison |
| `mind map` | Hierarchical concept map |
| `explain simply` | ELI15 explanation |
"""

USER_AVATAR = "https://em-content.zobj.net/source/twitter/376/bust-in-silhouette_1f464.png"
BOT_AVATAR  = "https://em-content.zobj.net/source/twitter/376/eagle_1f985.png"

# ==============================================================================
# 🏗️  BUILD UI
# ==============================================================================

# FIX 4: HTML-based password status — no textboxes
def _pw_html(locked: bool, msg: str = "") -> str:
    if locked:
        icon  = "🔒"
        cls   = "pw-locked"
        label = msg or "Locked — enter password"
    else:
        icon  = "✅"
        cls   = "pw-ok"
        label = msg or "Access granted"
    return f'<div class="pw-status-box {cls}">{icon} {label}</div>'


def apply_access(username: str, password_ok: bool):
    valid   = is_valid_username(username)
    allowed = valid and password_ok
    pw_html = _pw_html(not allowed, "Access granted" if allowed else "Locked — enter password")
    btn_cls = "unlock-btn-ok" if allowed else ""
    return (
        pw_html,
        btn_cls,
        gr.update(interactive=allowed),
        gr.update(interactive=allowed),
        gr.update(interactive=allowed),
        gr.update(interactive=allowed),
        gr.update(visible=password_ok),
        gr.update(visible=password_ok),
        gr.update(interactive=allowed),
        gr.update(interactive=allowed),
        gr.update(interactive=allowed),
        gr.update(interactive=allowed),
    )


def check_password(password: str, username: str):
    ok  = password.strip() == "187177"
    msg = "✅ Access granted" if ok else "❌ Wrong password — try again"
    pw_html = _pw_html(not ok, msg)
    btn_cls = "unlock-btn-ok" if ok else ""
    updates = apply_access(username, ok)
    # updates[0] is pw_html, updates[1] is btn_cls — replace with our values
    return (ok, pw_html, btn_cls, *updates[2:])


def on_username_change(username: str, password_ok: bool):
    return apply_access(username, password_ok)


def download_db_handler():
    path = get_db_path_for_download()
    if path and Path(path).exists():
        return gr.update(value=path, visible=True)
    return gr.update(value=None, visible=False)

def download_hist_handler(username: str):
    path = make_history_zip_for_user(username)
    if path and Path(path).exists():
        return gr.update(value=path, visible=True)
    return gr.update(value=None, visible=False)


def build_ui():
    with gr.Blocks(title="Phoenix Eduplan v3.0", css=DARK_CSS + NOTEBOOKLM_CSS) as demo:

        gr.HTML(HEADER_HTML)

        with gr.Row(equal_height=True):

            # ── Sidebar ──────────────────────────────────────────────────────
            with gr.Column(scale=3, min_width=300):

                # FIX 3: 5 visible tabs — merged Audio+Video into one Media tab,
                # all tabs fit without the "..." overflow button.
                with gr.Tabs():

                    # ── User tab ─────────────────────────────────────────────
                    with gr.Tab("👤 User"):
                        user_name = gr.Textbox(
                            label="Your username",
                            placeholder="Type a name to start",
                            lines=1,
                        )
                        # FIX 4: replaced two Textbox status fields with HTML badges
                        pw_status_html = gr.HTML(_pw_html(True))
                        pw_btn_wrapper = gr.HTML("")   # carries CSS class for glow

                        password     = gr.Textbox(
                            label="Password", type="password",
                            placeholder="Enter password…"
                        )
                        password_btn = gr.Button("🔓 Unlock", variant="primary")

                        download_db_btn   = gr.Button("⬇️ Download DB",           variant="secondary", visible=False)
                        download_hist_btn = gr.Button("⬇️ Download History ZIP",  variant="secondary", visible=False)
                        download_db_file  = gr.File(visible=False)
                        download_hist_file= gr.File(visible=False)

                    # ── Upload tab ───────────────────────────────────────────
                    with gr.Tab("📁 Upload"):
                        file_input = gr.File(
                            label="Drop files here",
                            file_types=[".pdf",".docx",".doc",".pptx",".ppt",
                                        ".xlsx",".xls",".csv",".txt",".md",
                                        ".png",".jpg",".jpeg",".webp"],
                            file_count="multiple",
                            type="filepath",
                            interactive=False,
                        )
                        process_btn = gr.Button("⚡ Index Files", variant="primary",  size="sm", interactive=False)
                        clear_btn   = gr.Button("🗑️ Clear Docs",  variant="stop",     size="sm", interactive=False)
                        status_box  = gr.Textbox(
                            label="Indexing Status",
                            value="📁 Ready — upload your files above",
                            lines=9, interactive=False,
                        )

                    # ── FIX 3: Combined Audio + Video tab ────────────────────
                    with gr.Tab("🎙️ Media"):
                        with gr.Tab("🎧 Audio"):
                            audio_text = gr.Textbox(
                                label="Text for narration (leave empty to use uploaded docs)",
                                lines=5,
                                placeholder="Type notes or leave empty to use uploaded documents…",
                                interactive=False,
                            )
                            audio_use_docs = gr.Checkbox(label="Use uploaded documents", value=True)
                            audio_lang     = gr.Dropdown(
                                ["Auto", "English", "Arabic"], value="Auto", label="Output language"
                            )
                            audio_btn    = gr.Button("🎧 Generate Audio", variant="primary", interactive=False)
                            audio_status = gr.Textbox(value="Ready.", interactive=False, show_label=False)
                            audio_out    = gr.Audio(label="Narration Audio")

                        with gr.Tab("🎬 Video"):
                            video_text = gr.Textbox(
                                label="Text for study video (leave empty to use uploaded docs)",
                                lines=5,
                                placeholder="Type notes or leave empty to use uploaded documents…",
                                interactive=False,
                            )
                            video_use_docs = gr.Checkbox(label="Use uploaded documents", value=True)
                            video_lang     = gr.Dropdown(
                                ["Auto", "English", "Arabic"], value="Auto", label="Output language"
                            )
                            video_btn    = gr.Button("🎬 Generate Video", variant="primary", interactive=False)
                            video_status = gr.Textbox(value="Ready.", interactive=False, show_label=False)
                            video_out    = gr.Video(label="Study Video")

                    # ── Docs tab ─────────────────────────────────────────────
                    with gr.Tab("📚 Docs"):
                        gr.Markdown("### Indexed Documents")
                        docs_display = gr.Textbox(
                            label="Active files",
                            value="No documents yet.",
                            lines=10, interactive=False,
                        )
                        refresh_btn = gr.Button("🔄 Refresh", variant="secondary", size="sm")

                    # ── Guide tab ────────────────────────────────────────────
                    with gr.Tab("💡 Guide"):
                        gr.Markdown(SIDEBAR_GUIDE)

            # ── Chat ─────────────────────────────────────────────────────────
            with gr.Column(scale=7):
                chat = gr.ChatInterface(
                    fn=chat_logic,
                    additional_inputs=[user_name],
                    chatbot=gr.Chatbot(
                        height=580,
                        avatar_images=(USER_AVATAR, BOT_AVATAR),
                        render_markdown=True,
                        placeholder=(
                            "<div style='text-align:center;padding:60px 20px;"
                            "font-family:DM Sans,sans-serif;position:relative;z-index:1'>"
                            "<div style='font-size:3rem'>🦅</div>"
                            "<div style='font-size:1.1rem;margin-top:12px;font-weight:500;"
                            "color:#d8c7bf'>Phoenix Eduplan v3.0</div>"
                            "<div style='font-size:0.85rem;margin-top:6px;color:#9e8880'>"
                            "Upload a document → then ask anything</div></div>"
                        ),
                    ),
                    textbox=gr.Textbox(
                        placeholder="Enter your name to start chatting…",
                        show_label=False,
                        scale=7,
                        interactive=False,
                    ),
                    examples=[
                        ["Summarize the document in structured bullet points"],
                        ["Generate a 5-question quiz with answers"],
                        ["Create flashcards for the key concepts"],
                        ["Draw a mind map of the main topics"],
                        ["Explain the main idea simply, like I'm 15"],
                        ["What are the most important terms and definitions?"],
                        ["Compare and contrast the key concepts"],
                        ["ما هي النقاط الرئيسية في هذا المستند؟"],
                        ["قارن بين المفاهيم الرئيسية في المستند"],
                    ],
                    cache_examples=False,
                )

        password_state = gr.State(False)

        # ── Wire up events ────────────────────────────────────────────────────

        # Shared outputs list (excluding the two pw outputs at the front)
        _shared_access_outputs = [
            chat.textbox, file_input, process_btn, clear_btn,
            download_db_btn, download_hist_btn,
            audio_text, audio_btn, video_text, video_btn,
        ]

        user_name.change(
            fn=on_username_change,
            inputs=[user_name, password_state],
            outputs=[pw_status_html, pw_btn_wrapper, *_shared_access_outputs],
        )

        password_btn.click(
            fn=check_password,
            inputs=[password, user_name],
            outputs=[password_state, pw_status_html, pw_btn_wrapper, *_shared_access_outputs],
        )

        process_btn.click(fn=process_files,   inputs=[file_input, user_name], outputs=status_box)
        file_input.upload(fn=process_files,   inputs=[file_input, user_name], outputs=status_box)
        clear_btn.click(  fn=clear_documents, inputs=user_name,               outputs=status_box)
        refresh_btn.click(fn=list_documents,  inputs=user_name,               outputs=docs_display)

        audio_btn.click(
            fn=generate_audio_only,
            inputs=[audio_text, user_name, audio_use_docs, audio_lang],
            outputs=[audio_out, audio_status],
        )
        video_btn.click(
            fn=generate_video_only,
            inputs=[video_text, user_name, video_use_docs, video_lang],
            outputs=[video_out, video_status],
        )

        download_db_btn.click(  fn=download_db_handler,              outputs=download_db_file)
        download_hist_btn.click(fn=download_hist_handler, inputs=user_name, outputs=download_hist_file)

    return demo

# ==============================================================================
# 🚀  LAUNCH + API
# ==============================================================================

demo = build_ui()

if FASTAPI_AVAILABLE:
    app = FastAPI()

    def require_internal_token(
        x_internal_token: Optional[str] = Header(None),
        authorization: Optional[str] = Header(None, alias="Authorization"),
    ):
        if not AI_SERVICE_TOKEN:
            return
        token = x_internal_token
        if not token and authorization:
            parts = authorization.split(" ", 1)
            if len(parts) == 2 and parts[0].lower() == "bearer":
                token = parts[1]
        if token != AI_SERVICE_TOKEN:
            raise HTTPException(status_code=403, detail="Invalid internal token.")

    @app.get("/api/health")
    def api_health(_=Depends(require_internal_token)):
        return {"status": "ok"}

    @app.post("/api/chat")
    def api_chat(payload: dict = Body(...), _=Depends(require_internal_token)):
        session_id = payload.get("session_id", "")
        username = payload.get("username", "")
        message  = payload.get("message", "")
        history  = payload.get("history", [])
        if not session_id:
            raise HTTPException(status_code=400, detail="session_id is required")
        final    = ""
        for chunk in chat_logic(message, history, username, session_id=session_id):
            final = chunk
        return {"reply": final, "session_id": session_id}

    @app.post("/api/upload")
    async def api_upload(
        files: list[UploadFile] = File(...),
        username: Optional[str] = Form(None),
        session_id: Optional[str] = Form(None),
        _=Depends(require_internal_token),
    ):
        if not session_id or not username:
            raise HTTPException(status_code=400, detail="session_id and username are required")
        paths = []
        for f in files:
            suffix = Path(f.filename).suffix
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(await f.read())
                paths.append(tmp.name)
        result = process_files(paths, username, session_id=session_id)
        for p in paths:
            try:
                os.remove(p)
            except Exception:
                pass
        return {"status": result}

    app = gr.mount_gradio_app(app, demo, path="/")

if __name__ == "__main__":
    demo.launch(
        server_name="0.0.0.0",
        server_port=7860,
        ssr_mode=False,
    )
