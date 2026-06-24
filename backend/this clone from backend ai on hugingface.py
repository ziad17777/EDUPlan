"""
Phoenix Eduplan v4.1
Advanced AI Study Assistant — HuggingFace Spaces
"""

# ==============================================================================
# 📦  IMPORTS
# ==============================================================================
# Loads all necessary external libraries, environment variables, and Hugging Face tokens.

import gradio as gr
import os
import re
import hashlib
import sqlite3
import json
import uuid
import shutil
import tempfile
import zipfile
import time
import logging
import threading
import math
import textwrap
from pathlib import Path
from datetime import datetime
from typing import Generator, Optional
from collections import defaultdict
from contextlib import contextmanager
from functools import lru_cache

# Core ML / RAG
from huggingface_hub import InferenceClient
from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, TextLoader, CSVLoader, UnstructuredMarkdownLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document

# ── Optional: pdfplumber ─────────────────────────────────────────────────────
PDFPLUMBER_AVAILABLE = False
try:
    import pdfplumber
    PDFPLUMBER_AVAILABLE = True
except ImportError:
    pass

# ── Optional: python-pptx ────────────────────────────────────────────────────
PPTX_AVAILABLE = False
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    pass

# ── Optional: pandas ─────────────────────────────────────────────────────────
PANDAS_AVAILABLE = False
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    pass

# ── Optional: PIL ─────────────────────────────────────────────────────────────
PIL_AVAILABLE = False
try:
    from PIL import Image, ImageDraw, ImageFont, ImageFilter
    PIL_AVAILABLE = True
except ImportError:
    pass

# ── Optional: pytesseract ─────────────────────────────────────────────────────
TESSERACT_AVAILABLE = False
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    pass

# ── Optional: BM25 ────────────────────────────────────────────────────────────
BM25_AVAILABLE = False
try:
    from rank_bm25 import BM25Okapi
    BM25_AVAILABLE = True
except ImportError:
    pass

# ── Optional: CrossEncoder (reranking) ────────────────────────────────────────
RERANKER_AVAILABLE = False
try:
    from sentence_transformers import CrossEncoder
    RERANKER_AVAILABLE = True
except ImportError:
    pass

# ── Optional: FastAPI ─────────────────────────────────────────────────────────
FASTAPI_AVAILABLE = False
try:
    from fastapi import FastAPI, UploadFile, File, Form
    from fastapi.responses import JSONResponse
    FASTAPI_AVAILABLE = True
except ImportError:
    pass

# ── Optional: Azure Speech ────────────────────────────────────────────────────
AZURE_TTS_AVAILABLE = False
try:
    import azure.cognitiveservices.speech as speechsdk
    AZURE_TTS_AVAILABLE = True
except ImportError:
    pass

# ── Optional: gTTS ────────────────────────────────────────────────────────────
GTTS_AVAILABLE = False
try:
    from gtts import gTTS
    GTTS_AVAILABLE = True
except ImportError:
    pass

# ── Optional: moviepy ─────────────────────────────────────────────────────────
MOVIEPY_AVAILABLE = False
try:
    from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
    MOVIEPY_AVAILABLE = True
except Exception:
    pass

# ==============================================================================
# ⚙️  CONFIGURATION
# ==============================================================================
# Defines global variables, model repo IDs, and UI constraints (like max tokens and UI dimensions).

REPO_ID      = "meta-llama/Llama-3.3-70B-Instruct"
EMBED_MODEL  = "BAAI/bge-m3"
RERANK_MODEL = "BAAI/bge-reranker-v2-m3"
IMAGE_MODEL  = "black-forest-labs/FLUX.1-schnell"
TTS_MODEL_EN = "facebook/mms-tts-eng"
TTS_MODEL_AR = "facebook/mms-tts-ara"

DATA_DIR     = Path("/data")
DB_PATH      = str(DATA_DIR / "phoenix_history.db")
FAISS_DIR    = DATA_DIR / "faiss_indexes"

MAX_HISTORY_TURNS               = 14
MAX_CONTEXT_CHARS               = 16000
SUMMARY_EVERY_N_ASSISTANT_TURNS = 5
MAX_RESPONSE_TOKENS             = 2500

RATE_LIMIT_WINDOW    = 60
RATE_LIMIT_MAX_CALLS = 15

HF_TOKEN = os.getenv("HF_TOKEN")
if not HF_TOKEN:
    raise ValueError("❌ HF_TOKEN missing!")

AZURE_SPEECH_KEY    = os.getenv("AZURE_SPEECH_KEY", "")
AZURE_SPEECH_REGION = os.getenv("AZURE_SPEECH_REGION", "francecentral")

client = InferenceClient(model=REPO_ID, token=HF_TOKEN)

# ── Logging ──────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("phoenix")

# ==============================================================================
# 🛠️  HELPERS
# ==============================================================================
# Contains utility functions for text extraction and instantiates global singleton ML models.

def _extract_text(content) -> str:
    """Safely extract text from various content types."""
    if content is None:
        return ""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
            elif isinstance(item, dict):
                parts.append(item.get("text", str(item)))
            else:
                parts.append(str(item))
        return "".join(parts)
    return str(content)

# ── Embeddings singleton ─────────────────────────────────────────────────────
_embed_model = None
_embed_lock  = threading.Lock()

def get_embeddings():
    global _embed_model
    if _embed_model is None:
        with _embed_lock:
            if _embed_model is None:
                log.info(f"Loading embedding model: {EMBED_MODEL}")
                _embed_model = HuggingFaceEmbeddings(
                    model_name=EMBED_MODEL,
                    model_kwargs={"device": "cpu"},
                    encode_kwargs={"normalize_embeddings": True},
                )
                log.info("Embedding model loaded.")
    return _embed_model

# ── Reranker singleton ────────────────────────────────────────────────────────
_reranker      = None
_reranker_lock = threading.Lock()

def get_reranker():
    global _reranker
    if not RERANKER_AVAILABLE:
        return None
    if _reranker is None:
        with _reranker_lock:
            if _reranker is None:
                try:
                    log.info(f"Loading reranker: {RERANK_MODEL}")
                    _reranker = CrossEncoder(RERANK_MODEL, max_length=512)
                    log.info("Reranker loaded.")
                except Exception as e:
                    log.error(f"Reranker load failed: {e}")
                    return None
    return _reranker

# ==============================================================================
# 🗄️  DATABASE HELPER
# ==============================================================================
# Manages the SQLite connection pool using a context manager with WAL concurrency.

@contextmanager
def get_db():
    """Context manager for safe SQLite connections. Uses WAL mode to prevent locking during concurrent reads/writes."""
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.execute("PRAGMA journal_mode=WAL;")
    conn.execute("PRAGMA synchronous=NORMAL;")
    try:
        yield conn
        conn.commit()
    except Exception:
        conn.rollback()
        raise
    finally:
        conn.close()

# ==============================================================================
# 🚦  RATE LIMITER
# ==============================================================================
# Prevents API spam by enforcing a strict sliding-window request limit per user.

_rate_store: dict = {}
_rate_lock = threading.Lock()

def rate_limit_ok(username: str) -> bool:
    now = time.time()
    with _rate_lock:
        window = _rate_store.setdefault(username, [])
        window[:] = [t for t in window if now - t < RATE_LIMIT_WINDOW]
        if len(window) >= RATE_LIMIT_MAX_CALLS:
            return False
        window.append(now)
    return True

# ==============================================================================
# 👤  PER-USER CONTEXTS
# ==============================================================================
# A thread-safe LRU cache that stores and evicts FAISS vector stores from RAM to prevent server crashes.

from collections import OrderedDict

class LRUUserCache:
    """A simple LRU cache for user contexts to prevent unbounded memory growth."""
    def __init__(self, capacity=50):
        self.cache = OrderedDict()
        self.capacity = capacity

    def get(self, key, default=None):
        if key not in self.cache:
            return default
        self.cache.move_to_end(key)
        return self.cache[key]

    def set(self, key, value):
        if key in self.cache:
            self.cache.move_to_end(key)
        self.cache[key] = value
        if len(self.cache) > self.capacity:
            self.cache.popitem(last=False)

    def __contains__(self, key):
        return key in self.cache

    def __getitem__(self, key):
        return self.get(key)
        
    def __setitem__(self, key, value):
        self.set(key, value)

_user_contexts = LRUUserCache(capacity=50)
_ctx_lock = threading.Lock()

def load_user_indexes(username: str, ctx: dict):
    idx_base = FAISS_DIR / sanitize_username(username)
    if not idx_base.exists() or not idx_base.is_dir():
        return
        
    embed = get_embeddings()
    merged_vs = None
    
    for doc_dir in idx_base.iterdir():
        if doc_dir.is_dir():
            doc_id = doc_dir.name
            try:
                vs = FAISS.load_local(str(doc_dir), embed, allow_dangerous_deserialization=True)
                if merged_vs is None:
                    merged_vs = vs
                else:
                    merged_vs.merge_from(vs)
                
                # Reconstruct doc_store for BM25
                chunks = list(vs.docstore._dict.values())
                if chunks:
                    ctx["doc_store"][doc_id] = chunks
                    ctx["active_docs"].add(doc_id)
                    ctx["doc_hashes"][doc_id] = chunks[0].metadata.get("file_hash", "reloaded_" + doc_id)
            except Exception as e:
                log.error(f"Failed to load FAISS for {doc_id}: {e}")
                
    if merged_vs:
        ctx["vector_store"] = merged_vs
        rebuild_bm25(ctx)
        log.info(f"Loaded persistent indexes for {username}")


def get_user_context(username: str) -> dict:
    """Manages multi-tenant server memory using a Thread-Safe LRU Cache. Automatically evicts dormant FAISS vectors to prevent out-of-memory (OOM) crashes."""
    with _ctx_lock:
        if username not in _user_contexts:
            _user_contexts[username] = {
                "vector_store": None,
                "doc_store":    {},
                "active_docs":  set(),
                "doc_hashes":   {},
                "bm25":         None,
                "bm25_docs":    [],
                "session_id":   None,
            }
            load_user_indexes(username, _user_contexts[username])
        return _user_contexts[username]

def clear_user_context(username: str):
    with _ctx_lock:
        _user_contexts[username] = {
            "vector_store": None,
            "doc_store":    {},
            "active_docs":  set(),
            "doc_hashes":   {},
            "bm25":         None,
            "bm25_docs":    [],
            "session_id":   _user_contexts.get(username, {}).get("session_id"),
        }
        user_faiss_dir = os.path.join(FAISS_DIR, sanitize_username(username))
        if os.path.exists(user_faiss_dir):
            import shutil
            shutil.rmtree(user_faiss_dir, ignore_errors=True)

def rebuild_bm25(ctx: dict):
    if not BM25_AVAILABLE:
        return
    corpus_docs = []
    for did, chunks in ctx["doc_store"].items():
        if did in ctx["active_docs"]:
            corpus_docs.extend(chunks)
    if corpus_docs:
        tokenized = [d.page_content.lower().split() for d in corpus_docs]
        ctx["bm25"] = BM25Okapi(tokenized)
        ctx["bm25_docs"] = corpus_docs
    else:
        ctx["bm25"] = None
        ctx["bm25_docs"] = []

# ==============================================================================
# 🔐  USERNAME + DATA DIRS
# ==============================================================================
# Handles user sanitization, directory generation, and file path assignments.

def sanitize_username(name: str) -> str:
    if not name:
        return ""
    return re.sub(r"[^a-zA-Z0-9_\-\u0600-\u06FF ]", "", name.strip())[:50]

def is_valid_username(name: str) -> bool:
    return len(sanitize_username(name)) >= 2

def ensure_user_dirs(username: str):
    safe = sanitize_username(username)
    (DATA_DIR / "users" / safe).mkdir(parents=True, exist_ok=True)
    (FAISS_DIR / safe).mkdir(parents=True, exist_ok=True)

# ==============================================================================
# 🗃️  DATABASE OPERATIONS
# ==============================================================================
# Executes CRUD operations for chat history and document session tracking via SQLite.

def init_db():
    DATA_DIR.mkdir(exist_ok=True)
    with get_db() as conn:
        c = conn.cursor()
        c.execute("""CREATE TABLE IF NOT EXISTS sessions (
            id TEXT PRIMARY KEY,
            username TEXT NOT NULL DEFAULT '',
            name TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            doc_names TEXT DEFAULT '[]',
            summary TEXT DEFAULT ''
        )""")
        c.execute("""CREATE TABLE IF NOT EXISTS messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_id TEXT NOT NULL,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            sources TEXT DEFAULT '[]',
            confidence REAL DEFAULT 0.0,
            created_at TEXT NOT NULL,
            FOREIGN KEY (session_id) REFERENCES sessions(id)
        )""")
        c.execute("CREATE INDEX IF NOT EXISTS idx_messages_session ON messages(session_id)")
        c.execute("CREATE INDEX IF NOT EXISTS idx_sessions_username ON sessions(username)")


def create_session(username="", name=None, doc_names=None):
    session_id = str(uuid.uuid4())[:8]
    now = datetime.utcnow().isoformat()
    if not name:
        name = f"Session {datetime.now().strftime('%b %d %H:%M')}"
    with get_db() as conn:
        conn.execute(
            "INSERT INTO sessions (id,username,name,created_at,updated_at,doc_names) VALUES (?,?,?,?,?,?)",
            (session_id, sanitize_username(username), name, now, now, json.dumps(doc_names or []))
        )
    return session_id


def save_message(session_id, role, content, sources=None, confidence=0.0):
    now = datetime.utcnow().isoformat()
    with get_db() as conn:
        conn.execute(
            "INSERT INTO messages (session_id,role,content,sources,confidence,created_at) VALUES (?,?,?,?,?,?)",
            (session_id, role, content, json.dumps(sources or []), confidence, now)
        )
        conn.execute("UPDATE sessions SET updated_at=? WHERE id=?", (now, session_id))


def get_session_summary(session_id):
    with get_db() as conn:
        row = conn.execute("SELECT summary FROM sessions WHERE id=?", (session_id,)).fetchone()
    return row[0] if row else ""


def update_session_summary(session_id, summary):
    with get_db() as conn:
        conn.execute(
            "UPDATE sessions SET summary=?, updated_at=? WHERE id=?",
            (summary, datetime.utcnow().isoformat(), session_id)
        )


def load_session_messages(session_id, limit=40):
    with get_db() as conn:
        rows = conn.execute(
            "SELECT role, content, sources, confidence, created_at FROM messages "
            "WHERE session_id=? ORDER BY id DESC LIMIT ?",
            (session_id, limit)
        ).fetchall()
    rows.reverse()
    return [{"role": r[0], "content": r[1], "sources": json.loads(r[2]),
             "confidence": r[3], "created_at": r[4]} for r in rows]


def count_assistant_turns(session_id):
    with get_db() as conn:
        row = conn.execute(
            "SELECT COUNT(*) FROM messages WHERE session_id=? AND role='assistant'",
            (session_id,)
        ).fetchone()
    return row[0] if row else 0


def update_session_docs(session_id, doc_names):
    with get_db() as conn:
        conn.execute(
            "UPDATE sessions SET doc_names=?, updated_at=? WHERE id=?",
            (json.dumps(doc_names), datetime.utcnow().isoformat(), session_id)
        )


def list_user_sessions(username, limit=20):
    with get_db() as conn:
        rows = conn.execute(
            "SELECT id, name, updated_at, doc_names FROM sessions "
            "WHERE username=? ORDER BY updated_at DESC LIMIT ?",
            (sanitize_username(username), limit)
        ).fetchall()
    return [{"id": r[0], "name": r[1], "updated_at": r[2],
             "doc_names": json.loads(r[3])} for r in rows]


init_db()


def get_or_create_user_session(username):
    ctx = get_user_context(username)
    if ctx["session_id"]:
        return ctx["session_id"]
    sessions = list_user_sessions(username, limit=1)
    if sessions:
        sid = sessions[0]["id"]
        ctx["session_id"] = sid
        return sid
    sid = create_session(username=username)
    ctx["session_id"] = sid
    return sid


def switch_user_session(username, session_id):
    ctx = get_user_context(username)
    ctx["session_id"] = session_id
    msgs = load_session_messages(session_id)
    history = []
    for m in msgs:
        if m["role"] == "user":
            history.append({"role": "user", "content": m["content"]})
        elif m["role"] == "assistant":
            history.append({"role": "assistant", "content": m["content"]})
    return history

# ==============================================================================
# 📂  USER HISTORY FILES
# ==============================================================================
# Manages exporting and importing raw JSON user chat logs.

def save_user_history(username, history):
    if not is_valid_username(username):
        return
    try:
        ensure_user_dirs(username)
        safe = sanitize_username(username)
        path = DATA_DIR / "users" / safe / "chat_history.json"
        with open(path, "w", encoding="utf-8") as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
    except Exception as e:
        log.error(f"save_user_history failed for {username}: {e}")


def load_user_history(username):
    if not is_valid_username(username):
        return []
    try:
        safe = sanitize_username(username)
        path = DATA_DIR / "users" / safe / "chat_history.json"
        if path.exists():
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception as e:
        log.error(f"load_user_history failed for {username}: {e}")
    return []


def download_db():
    if os.path.exists(DB_PATH):
        # Create temp file in the current working directory to allow Gradio to serve it
        tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".db", delete=False)
        shutil.copy2(DB_PATH, tmp.name)
        return tmp.name
    return None


def download_user_history(username):
    if not is_valid_username(username):
        return None
    try:
        safe = sanitize_username(username)
        user_dir = DATA_DIR / "users" / safe
        if not user_dir.exists():
            return None
        # Create temp file in the current working directory to allow Gradio to serve it
        tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".zip", delete=False)
        with zipfile.ZipFile(tmp.name, "w", zipfile.ZIP_DEFLATED) as zf:
            for fp in user_dir.rglob("*"):
                if fp.is_file():
                    zf.write(fp, fp.relative_to(user_dir))
        return tmp.name
    except Exception as e:
        log.error(f"download_user_history failed for {username}: {e}")
        return None

# ==============================================================================
# 📄  FILE LOADERS
# ==============================================================================
# Parses PDFs, TXTs, and DOCX files using LangChain loaders with automatic OCR fallback.

def load_pdf(path):
    """Load PDF with pdfplumber fallback."""
    docs = []
    try:
        loader = PyPDFLoader(path)
        docs = loader.load()
    except Exception:
        pass

    if not docs and PDFPLUMBER_AVAILABLE:
        try:
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    tables = page.extract_tables()
                    for table in tables:
                        for row in table:
                            text += "\n" + " | ".join([str(c) if c else "" for c in row])
                    if text.strip():
                        docs.append(Document(
                            page_content=text.strip(),
                            metadata={"source": os.path.basename(path), "page": i + 1}
                        ))
        except Exception:
            pass

    if not docs:
        try:
            loader = TextLoader(path, encoding="utf-8")
            docs = loader.load()
        except Exception:
            pass

    return docs


def load_docx(path):
    """Load DOCX files."""
    try:
        loader = Docx2txtLoader(path)
        return loader.load()
    except Exception:
        return []


def load_pptx(path):
    """Load PPTX files."""
    if not PPTX_AVAILABLE:
        return []
    docs = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                docs.append(Document(
                    page_content="\n".join(texts),
                    metadata={"source": os.path.basename(path), "slide": i + 1}
                ))
    except Exception:
        pass
    return docs


def load_excel(path):
    """Load Excel files."""
    if not PANDAS_AVAILABLE:
        return []
    docs = []
    try:
        xls = pd.ExcelFile(path)
        for sheet in xls.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet)
            text = f"Sheet: {sheet}\n{df.to_string(index=False)}"
            docs.append(Document(
                page_content=text,
                metadata={"source": os.path.basename(path), "sheet": sheet}
            ))
    except Exception:
        pass
    return docs


def load_csv(path):
    """Load CSV files."""
    try:
        loader = CSVLoader(path, encoding="utf-8")
        return loader.load()
    except Exception:
        try:
            loader = CSVLoader(path, encoding="latin-1")
            return loader.load()
        except Exception:
            return []


def load_image(path):
    """Load image files with OCR if available."""
    if not PIL_AVAILABLE:
        return []
    docs = []
    try:
        img = Image.open(path)
        text = ""
        if TESSERACT_AVAILABLE:
            try:
                text = pytesseract.image_to_string(img, lang="ara+eng")
            except Exception:
                pass
        if not text.strip():
            text = f"[Image file: {os.path.basename(path)}, size: {img.size}, mode: {img.mode}]"
        docs.append(Document(
            page_content=text.strip(),
            metadata={"source": os.path.basename(path), "type": "image"}
        ))
    except Exception:
        pass
    return docs


def load_file(path):
    """Route file to appropriate loader based on extension."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return load_pdf(path)
    elif ext in (".docx", ".doc"):
        return load_docx(path)
    elif ext in (".pptx", ".ppt"):
        return load_pptx(path)
    elif ext in (".xlsx", ".xls"):
        return load_excel(path)
    elif ext == ".csv":
        return load_csv(path)
    elif ext == ".md":
        try:
            loader = UnstructuredMarkdownLoader(path)
            return loader.load()
        except Exception:
            try:
                loader = TextLoader(path, encoding="utf-8")
                return loader.load()
            except Exception:
                return []
    elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
        return load_image(path)
    elif ext in (".txt", ".log", ".json", ".xml", ".html", ".htm", ".yaml", ".yml"):
        try:
            loader = TextLoader(path, encoding="utf-8")
            return loader.load()
        except Exception:
            return []
    else:
        try:
            loader = TextLoader(path, encoding="utf-8")
            return loader.load()
        except Exception:
            return []

# ==============================================================================
# 🧹  TEXT CLEANING
# ==============================================================================
# Normalizes messy text data and strips zero-width Unicode characters before embedding.

def clean_text(text: str) -> str:
    """Clean extracted text."""
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = re.sub(r"[^\S\n]+", " ", text)
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        stripped = line.strip()
        if stripped:
            cleaned.append(stripped)
        elif cleaned and cleaned[-1] != "":
            cleaned.append("")
    return "\n".join(cleaned).strip()

# ==============================================================================
# 📊  INDEXING
# ==============================================================================
# Chunks documents and builds dual dense/sparse retrieval indices (FAISS & BM25).

def build_index(docs, username="", doc_id=""):
    """
    Parses and indexes loaded documents using a RecursiveCharacterTextSplitter (800 char chunk, 200 overlap). 
    Builds both a dense FAISS vector space for semantic similarity and a sparse BM25 index for exact keyword matching.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800, chunk_overlap=200,   # ⬆️ Improved from 600/120
        separators=["\n\n", "\n", ". ", "! ", "? ", "، ", "؛ ", " ", ""],
    )
    chunks = splitter.split_documents(docs)

    if not chunks:
        return None

    for chunk in chunks:
        chunk.page_content = clean_text(chunk.page_content)
        chunk.metadata["doc_id"] = doc_id

    chunks = [c for c in chunks if len(c.page_content.strip()) > 20]

    if not chunks:
        return None

    embed = get_embeddings()
    vs = FAISS.from_documents(chunks, embed)

    if username:
        try:
            idx_path = FAISS_DIR / sanitize_username(username) / doc_id
            vs.save_local(str(idx_path))
        except Exception as e:
            log.error(f"Failed to save FAISS index for {doc_id}: {e}")

    return vs


def merge_indexes(base_vs, new_vs):
    """Merge two FAISS vector stores."""
    if base_vs is None:
        return new_vs
    if new_vs is None:
        return base_vs
    base_vs.merge_from(new_vs)
    return base_vs

# ==============================================================================
# 🔍  HYBRID RETRIEVAL
# ==============================================================================
# Executes Vector + Keyword search, mathematically merging them via Reciprocal Rank Fusion.

def hybrid_search(query, username, top_k=6):
    """
    Executes a dual-pipeline retrieval mechanism by running Vector Similarity Search and Lexical Search (BM25) simultaneously.
    It normalizes their distinct mathematical scoring models using a Reciprocal Rank Fusion (RRF) algorithm, and optionally passes the best chunks through a CrossEncoder reranker.
    """
    """Combine vector similarity with BM25 keyword search + optional reranking."""
    ctx        = get_user_context(username)
    vec_store  = ctx["vector_store"]
    doc_store  = ctx["doc_store"]
    active     = ctx["active_docs"]

    if not vec_store and not doc_store:
        return []

    all_results = []

    vec_results = []
    bm25_results = []

    # ── Vector search ─────────────────────────────────────────────────────────
    if vec_store:
        try:
            vec_docs = vec_store.similarity_search(query, k=top_k, fetch_k=max(20, top_k * 4), filter=lambda md: md.get("doc_id", "") in active)
            vec_results = list(vec_docs)
            all_results.extend(vec_results)
        except Exception as e:
            log.error(f"Vector search failed: {e}")

    # ── BM25 search ───────────────────────────────────────────────────────────
    if BM25_AVAILABLE and ctx.get("bm25") and ctx.get("bm25_docs"):
        try:
            bm25 = ctx["bm25"]
            corpus_docs = ctx["bm25_docs"]
            scores = bm25.get_scores(query.lower().split())
            ranked_idx = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)
            for i in ranked_idx[:top_k * 2]:
                if scores[i] > 0:
                    bm25_results.append(corpus_docs[i])
            all_results.extend(bm25_results)
        except Exception as e:
            log.error(f"BM25 search failed: {e}")

    if not all_results:
        return []

    # ── Reciprocal Rank Fusion (RRF) dedup ────────────────────────────────────
    seen = {}
    vec_rank = {}
    bm25_rank = {}

    for rank, doc in enumerate(vec_results):
        key = hashlib.md5(doc.page_content.encode()).hexdigest()
        if key not in seen:
            seen[key] = doc
        if key not in vec_rank:
            vec_rank[key] = rank
            
    for rank, doc in enumerate(bm25_results):
        key = hashlib.md5(doc.page_content.encode()).hexdigest()
        if key not in seen:
            seen[key] = doc
        if key not in bm25_rank:
            bm25_rank[key] = rank

    rrf_scores = {}
    for key in seen:
        rrf = 0.0
        if key in vec_rank:
            rrf += 1.0 / (60 + vec_rank[key])
        if key in bm25_rank:
            rrf += 1.0 / (60 + bm25_rank[key])
        rrf_scores[key] = rrf

    sorted_keys = sorted(rrf_scores, key=rrf_scores.get, reverse=True)
    merged = [seen[k] for k in sorted_keys[:top_k * 2]]

    # ── Reranking ─────────────────────────────────────────────────────────────
    reranker = get_reranker()
    if reranker and len(merged) > 1:
        try:
            pairs = [(query, d.page_content) for d in merged]
            scores = reranker.predict(pairs)
            ranked = sorted(zip(scores, merged), key=lambda x: x[0], reverse=True)
            merged = [doc for _, doc in ranked[:top_k]]
        except Exception as e:
            log.error(f"Reranking failed: {e}")
            merged = merged[:top_k]
    else:
        merged = merged[:top_k]

    return merged

# ==============================================================================
# 📥  DOCUMENT PROCESSING
# ==============================================================================
# Orchestrates user uploads, caching, vectorization, and UI status updates.

def process_files(files, username):
    """
    Orchestrates the ingestion of user uploads (PDF, TXT, DOCX) into the multi-tenant LRU cache.
    Automatically chunks the documents, creates their vector embeddings, and generates a short AI summary of the file to confirm successful processing.
    """
    if not is_valid_username(username):
        return "⚠️ Please enter your name first."
    if not files:
        return "📁 Ready — upload one or more files to get started."

    ensure_user_dirs(username)
    # ✅ REMOVED: clear_user_context(username) — was wiping previous uploads
    ctx         = get_user_context(username)
    doc_store   = ctx["doc_store"]
    active_docs = ctx["active_docs"]
    doc_hashes  = ctx["doc_hashes"]

    session_id = get_or_create_user_session(username)
    results    = []
    new_doc_names = []

    for f in files:
        path = f.name if hasattr(f, "name") else str(f)
        fname = os.path.basename(path)

        try:
            with open(path, "rb") as fh:
                file_hash = hashlib.md5(fh.read()).hexdigest()
        except Exception:
            results.append(f"❌ Cannot read: {fname}")
            continue

        if file_hash in doc_hashes.values():
            results.append(f"⚠️ Already loaded: {fname}")
            continue

        docs = load_file(path)
        if not docs:
            results.append(f"❌ No text extracted: {fname}")
            continue

        doc_id = f"{sanitize_username(username)}_{hashlib.md5(fname.encode()).hexdigest()[:8]}"
        for d in docs:
            d.metadata["doc_id"]  = doc_id
            d.metadata["source"]  = fname
            d.metadata["file_hash"] = file_hash

        vs = build_index(docs, username, doc_id)
        if vs:
            ctx["vector_store"] = merge_indexes(ctx["vector_store"], vs)

        doc_store[doc_id]   = docs
        active_docs.add(doc_id)
        doc_hashes[doc_id]  = file_hash
        new_doc_names.append(fname)

        n_chunks = len(docs)
        total_chars = sum(len(d.page_content) for d in docs)
        
        # Auto-summarize the first few chunks
        doc_text = " ".join([d.page_content for d in docs[:5]])
        doc_lang = detect_language(doc_text)
        summary_text = auto_summarize_document(doc_text, lang=doc_lang)
        summary_msg = f"\n   ↳ 📝 _Summary: {summary_text}_" if summary_text else ""
        
        results.append(f"✅ {fname} — {n_chunks} sections, {total_chars:,} chars{summary_msg}")

    if new_doc_names:
        rebuild_bm25(ctx)
        try:
            update_session_docs(session_id, list(active_docs))
        except Exception as e:
            log.error(f"Failed to update session docs: {e}")

    summary = "\n".join(results)
    total_docs = len(active_docs)
    return f"📊 **{total_docs} document(s) active**\n\n{summary}"


def clear_documents(username):
    if not is_valid_username(username):
        return "⚠️ Enter your name first.", []
    clear_user_context(username)
    return "🗑️ All documents cleared.", []


def list_documents(username):
    if not is_valid_username(username):
        return "⚠️ Enter your name first."
    ctx = get_user_context(username)
    doc_store   = ctx["doc_store"]
    active_docs = ctx["active_docs"]

    if not doc_store:
        return "📁 No documents loaded."

    lines = []
    for did in doc_store:
        status = "✅" if did in active_docs else "❌"
        n_chunks = len(doc_store[did])
        source = doc_store[did][0].metadata.get("source", did) if doc_store[did] else did
        lines.append(f"{status} **{source}** — {n_chunks} sections (ID: `{did}`)")
    return "\n".join(lines)


def toggle_document(username, doc_id, enable=True):
    if not is_valid_username(username):
        return "⚠️ Enter your name first."
    ctx = get_user_context(username)
    if doc_id not in ctx["doc_store"]:
        return f"❌ Document `{doc_id}` not found."
    if enable:
        ctx["active_docs"].add(doc_id)
        rebuild_bm25(ctx)
        return f"✅ Enabled `{doc_id}`"
    else:
        ctx["active_docs"].discard(doc_id)
        rebuild_bm25(ctx)
        return f"❌ Disabled `{doc_id}`"


def get_enabled_docs(username):
    if not is_valid_username(username):
        return []
    ctx = get_user_context(username)
    names = []
    for did in ctx["active_docs"]:
        chunks = ctx["doc_store"].get(did, [])
        if chunks:
            names.append(chunks[0].metadata.get("source", did))
        else:
            names.append(did)
    return names

# ==============================================================================
# 🌍  LANGUAGE DETECTION + FORMATTING
# ==============================================================================
# Automatically detects Arabic/English scripts and applies Right-to-Left formatting.

DOC_REQUEST_PATTERNS = [
    r"(من|في|عن|ما|هل|اشرح|لخص|وضح|اذكر|عرف|قارن|حلل|ناقش|صف|بين).*"
    r"(المحاضر|الملف|المستند|الوثيق|الكتاب|النص|المقال|الدرس|المادة|الفصل|الباب)",
    r"(summarize|explain|describe|analyze|compare|define|discuss|list|outline|what|how|why|from).*"
    r"(document|file|lecture|text|book|chapter|article|material|slide|pdf|notes)",
    r"(لخص|اشرح|وضح|حلل|ناقش|صف|بين|اذكر)\s+(لي|لنا)?\s*(ما|ماذا|كيف|لماذا|أين|متى)?",
]

def detect_language(text: str) -> str:
    arabic_chars = len(re.findall(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF]", text))
    total = len(text.strip())
    if total == 0:
        return "en"
    return "ar" if arabic_chars / total > 0.3 else "en"


def requires_document(text: str) -> bool:
    for pat in DOC_REQUEST_PATTERNS:
        if re.search(pat, text, re.IGNORECASE):
            return True
    return False


def to_arabic_numerals(text: str) -> str:
    western = "0123456789"
    eastern = "٠١٢٣٤٥٦٧٨٩"
    table = str.maketrans(western, eastern)
    return text.translate(table)


def format_arabic_response(text: str) -> str:
    """Apply Arabic formatting: RTL markers, numeral conversion, list fix."""
    if detect_language(text) != "ar":
        return text

    lines = text.split("\n")
    formatted = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            formatted.append("")
            continue

        if re.match(r"^\d+[\.\)]\s", stripped):
            num_match = re.match(r"^(\d+)([\.\)])\s(.*)$", stripped)
            if num_match:
                num = to_arabic_numerals(num_match.group(1))
                sep = num_match.group(2)
                rest = num_match.group(3)
                stripped = f"{num}{sep} {rest}"

        if re.match(r"^[-*•]\s", stripped):
            stripped = "• " + stripped[2:]

        if re.match(r"^#{1,6}\s", stripped):
            formatted.append(stripped)
        else:
            formatted.append(stripped)

    return "\n".join(formatted)


def export_session_markdown(username):
    """Export current session as markdown."""
    if not is_valid_username(username):
        return None
    ctx = get_user_context(username)
    sid = ctx.get("session_id")
    if not sid:
        return None

    msgs = load_session_messages(sid)
    if not msgs:
        return None

    lines = [f"# Phoenix Eduplan — Session Export", f"**User:** {username}", f"**Date:** {datetime.now().isoformat()}", ""]
    for m in msgs:
        role = "👤 User" if m["role"] == "user" else "🤖 Assistant"
        lines.append(f"## {role}")
        lines.append(m["content"])
        lines.append("")

    tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".md", delete=False, mode="w", encoding="utf-8")
    tmp.write("\n".join(lines))
    tmp.close()
    return tmp.name


def export_chat_to_file(history, username):
    """Export chat history to a downloadable file."""
    if not history:
        return None
    lines = [f"# Phoenix Eduplan Chat Export", f"**User:** {username}", f"**Date:** {datetime.now().isoformat()}", ""]
    for msg in history:
        role_label = "👤 User" if msg.get("role") == "user" else "🤖 Assistant"
        content = msg.get("content", "")
        lines.append(f"## {role_label}")
        lines.append(content)
        lines.append("")
    tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".md", delete=False, mode="w", encoding="utf-8")
    tmp.write("\n".join(lines))
    tmp.close()
    return tmp.name

# ==============================================================================
# ✏️  QUERY REWRITING
# ==============================================================================
# Uses Llama-3 to expand user queries for better semantic retrieval coverage.

@lru_cache(maxsize=128)
def rewrite_query(question: str, lang: str) -> str:
    """Rewrite short/vague queries for better retrieval."""
    if len(question.split()) >= 6:
        return question
    try:
        prompt = (
            f"Rewrite this short question into a detailed search query (same language, "
            f"no explanation, just the rewritten query):\n{question}"
        )
        resp = client.text_generation(prompt, max_new_tokens=60, temperature=0.3)
        rewritten = resp.strip()
        if rewritten and len(rewritten) > len(question):
            return rewritten
    except Exception:
        pass
    return question


def is_followup(message: str, history: list) -> bool:
    """Detect if message is a follow-up to previous conversation."""
    if not history:
        return False
    followup_markers = [
        "ماذا عن", "وماذا", "وكيف", "بالإضافة", "أيضا", "كذلك", "هل يمكن",
        "what about", "and how", "also", "additionally", "can you", "tell me more",
        "elaborate", "expand", "continue", "go on", "more details",
    ]
    lower_msg = message.lower().strip()
    if any(lower_msg.startswith(m) for m in followup_markers):
        return True
    pronouns = ["it", "this", "that", "these", "those", "هذا", "هذه", "ذلك", "تلك", "هو", "هي"]
    if any(p in lower_msg.split() for p in pronouns):
        return True
    return False


def build_contextual_query(message: str, history: list) -> str:
    """Build a contextual query incorporating conversation context."""
    if not is_followup(message, history):
        return message
    recent = history[-4:]
    context_parts = []
    for m in recent:
        if m.get("role") == "user":
            txt = _extract_text(m.get("content", ""))
            context_parts.append(txt)
    if context_parts:
        context = " | ".join(context_parts[-2:])
        return f"{context} → {message}"
    return message


def update_summary_if_needed(session_id, username):
    """Periodically generate a conversation summary."""
    try:
        n = count_assistant_turns(session_id)
        if n > 0 and n % SUMMARY_EVERY_N_ASSISTANT_TURNS == 0:
            msgs = load_session_messages(session_id, limit=20)
            text = "\n".join(
                f"{m['role']}: {m['content'][:200]}" for m in msgs[-10:]
            )
            prompt = (
                "Summarize this conversation in 2-3 sentences, capturing key topics "
                "and conclusions. Be concise:\n\n" + text
            )
            summary = client.text_generation(prompt, max_new_tokens=150, temperature=0.3)
            if summary:
                update_session_summary(session_id, summary.strip())
    except Exception as e:
        log.warning(f"Summary update failed for session {session_id}: {e}")

# ==============================================================================
# 🎯  MODE DETECTION + SYSTEM PROMPTS
# ==============================================================================
# Dynamically alters the LLM system instructions based on chat vs. study modes.

QUIZ_TRIGGERS     = ["quiz", "اختبار", "امتحان", "أسئلة", "اختبرني", "quiz me", "test me"]
FLASH_TRIGGERS    = ["flashcard", "بطاقات", "فلاش كارد", "flash card", "بطاقة"]
SUMMARY_TRIGGERS  = ["summarize", "لخص", "ملخص", "تلخيص", "summary", "sum up"]
MINDMAP_TRIGGERS  = ["mind map", "خريطة ذهنية", "خريطة مفاهيم", "concept map", "mindmap"]
COMPARE_TRIGGERS  = ["compare", "قارن", "مقارنة", "فرق بين", "difference between", "vs"]
EXPLAIN_TRIGGERS  = ["explain simply", "اشرح ببساطة", "explain like", "eli5", "eli15", "بسط"]
DEEPDIVE_TRIGGERS = ["deep dive", "تعمق", "in depth", "بالتفصيل", "detailed analysis", "تحليل مفصل"]
TEACH_TRIGGERS    = ["teach me", "علمني", "help me understand", "ساعدني أفهم", "tutor", "درسني"]
PRACTICE_TRIGGERS = ["practice", "تمارين", "problems", "مسائل", "exercises", "تدريبات", "practice problems"]
GUIDE_TRIGGERS    = ["study guide", "دليل دراسي", "دليل", "exam prep", "تحضير للامتحان", "revision"]
DEBATE_TRIGGERS   = ["debate", "ناقش", "both sides", "argue", "pros and cons", "إيجابيات وسلبيات"]
TIMELINE_TRIGGERS = ["timeline", "تسلسل", "chronological", "تاريخ", "history of", "تسلسل زمني"]


def detect_mode(message: str) -> str:
    """Detect the response mode based on user message."""
    lower = message.lower().strip()
    if any(t in lower for t in QUIZ_TRIGGERS):
        return "quiz"
    if any(t in lower for t in FLASH_TRIGGERS):
        return "flashcard"
    if any(t in lower for t in SUMMARY_TRIGGERS):
        return "summary"
    if any(t in lower for t in MINDMAP_TRIGGERS):
        return "mindmap"
    if any(t in lower for t in COMPARE_TRIGGERS):
        return "compare"
    if any(t in lower for t in EXPLAIN_TRIGGERS):
        return "explain"
    if any(t in lower for t in DEEPDIVE_TRIGGERS):
        return "deepdive"
    if any(t in lower for t in TEACH_TRIGGERS):
        return "teach"
    if any(t in lower for t in PRACTICE_TRIGGERS):
        return "practice"
    if any(t in lower for t in GUIDE_TRIGGERS):
        return "studyguide"
    if any(t in lower for t in DEBATE_TRIGGERS):
        return "debate"
    if any(t in lower for t in TIMELINE_TRIGGERS):
        return "timeline"
    return "chat"


def build_system_prompt(mode: str, lang: str, has_docs: bool, context: str = "") -> str:
    """Build system prompt based on mode and language."""
    base = (
        "You are Phoenix Eduplan, an advanced AI study assistant. "
        "You help students learn effectively through clear explanations, "
        "structured responses, and engaging study materials.\n\n"
    )

    if lang == "ar":
        base = (
            "أنت فينيكس إديوبلان، مساعد دراسي ذكي متقدم تتحدث باللهجة المصرية المحببة والواضحة. "
            "تساعد الطلاب على التعلم بفعالية من خلال شروحات واضحة وممتعة باللهجة المصرية "
            "وإجابات منظمة ومواد دراسية تفاعلية.\n\n"
        )

    if has_docs and context:
        if lang == "ar":
            base += f"استخدم السياق التالي من مستندات الطالب للإجابة:\n\n{context}\n\n"
            base += "إذا لم يكن السياق كافياً، أجب من معرفتك العامة مع التوضيح.\n"
        else:
            base += f"Use the following context from the student's documents:\n\n{context}\n\n"
            base += "If context is insufficient, answer from general knowledge and note this.\n"

    mode_prompts = {
        "quiz": (
            "Generate exactly 5 multiple-choice questions based on the material. "
            "Format each with A, B, C, D options. Put correct answers at the end."
            if lang == "en" else
            "أنشئ ٥ أسئلة اختيار من متعدد بناءً على المادة. "
            "لكل سؤال ٤ خيارات (أ، ب، ج، د). ضع الإجابات الصحيحة في النهاية."
        ),
        "flashcard": (
            "Create 6 study flashcards. Format each as:\n"
            "**Card N:**\n📝 Front: [question/term]\n💡 Back: [answer/definition]"
            if lang == "en" else
            "أنشئ ٦ بطاقات دراسية. لكل بطاقة:\n"
            "**بطاقة N:**\n📝 الوجه: [سؤال/مصطلح]\n💡 الخلف: [إجابة/تعريف]"
        ),
        "summary": (
            "Provide a structured summary with:\n"
            "1. Key concepts\n2. Main points\n3. Important details\n4. Conclusion"
            if lang == "en" else
            "قدم ملخصاً منظماً يتضمن:\n"
            "١. المفاهيم الرئيسية\n٢. النقاط الأساسية\n٣. التفاصيل المهمة\n٤. الخلاصة"
        ),
        "mindmap": (
            "Create a text-based mind map / concept hierarchy. Use indentation and symbols:\n"
            "🎯 Main Topic\n  ├── Subtopic 1\n  │   ├── Detail\n  │   └── Detail\n  └── Subtopic 2"
            if lang == "en" else
            "أنشئ خريطة ذهنية نصية. استخدم المسافات والرموز:\n"
            "🎯 الموضوع الرئيسي\n  ├── فرع ١\n  │   ├── تفصيل\n  │   └── تفصيل\n  └── فرع ٢"
        ),
        "compare": (
            "Create a detailed comparison using a structured format:\n"
            "| Aspect | Item 1 | Item 2 |\n|--------|--------|--------|\n"
            "Include similarities and differences."
            if lang == "en" else
            "أنشئ مقارنة تفصيلية بتنسيق منظم:\n"
            "| الجانب | العنصر ١ | العنصر ٢ |\n|--------|---------|----------|\n"
            "اذكر أوجه التشابه والاختلاف."
        ),
        "explain": (
            "Explain this concept simply, as if to a 15-year-old student. "
            "Use analogies, examples, and break complex ideas into simple steps."
            if lang == "en" else
            "اشرح هذا المفهوم ببساطة، كأنك تشرح لطالب في المرحلة الثانوية. "
            "استخدم التشبيهات والأمثلة وقسم الأفكار المعقدة إلى خطوات بسيطة."
        ),
        "deepdive": (
            "Provide an exhaustive, multi-section deep dive analysis. Cover: background, "
            "core concepts, mechanisms, applications, controversies, and future directions. "
            "Use headers, bullet points, and examples extensively. Aim for 800+ words."
            if lang == "en" else
            "قدم تحليلاً عميقاً وشاملاً متعدد الأقسام. غطِّ: الخلفية، المفاهيم الأساسية، "
            "الآليات، التطبيقات، الجدل، والاتجاهات المستقبلية. استخدم العناوين والنقاط والأمثلة بكثرة."
        ),
        "teach": (
            "Act as a Socratic tutor. Explain the concept step by step, then ask the student "
            "a thought-provoking question to check understanding. Use analogies from everyday life. "
            "End with a mini-challenge or reflection question."
            if lang == "en" else
            "تصرف كمعلم سقراطي. اشرح المفهوم خطوة بخطوة، ثم اسأل الطالب سؤالاً "
            "محفزاً للتفكير للتحقق من الفهم. استخدم تشبيهات من الحياة اليومية. "
            "اختم بتحدٍ صغير أو سؤال تأملي."
        ),
        "practice": (
            "Generate 10 practice problems of varying difficulty (easy → medium → hard). "
            "Include: 3 easy, 4 medium, 3 hard problems. Provide detailed solutions at the end. "
            "Label each problem with its difficulty level."
            if lang == "en" else
            "أنشئ ١٠ تمارين بمستويات صعوبة متدرجة (سهل ← متوسط ← صعب). "
            "تتضمن: ٣ سهلة، ٤ متوسطة، ٣ صعبة. قدم حلولاً مفصلة في النهاية. "
            "حدد مستوى الصعوبة لكل تمرين."
        ),
        "studyguide": (
            "Create a comprehensive study guide with:\n"
            "📌 Key concepts to master\n🎯 Learning objectives\n"
            "⚠️ Common mistakes to avoid\n💡 Exam tips\n"
            "📝 Quick-review checklist\n🔗 How topics connect to each other"
            if lang == "en" else
            "أنشئ دليل دراسة شامل يتضمن:\n"
            "📌 المفاهيم الرئيسية التي يجب إتقانها\n🎯 أهداف التعلم\n"
            "⚠️ الأخطاء الشائعة التي يجب تجنبها\n💡 نصائح للامتحان\n"
            "📝 قائمة مراجعة سريعة\n🔗 كيف ترتبط المواضيع ببعضها"
        ),
        "debate": (
            "Present a balanced debate analysis. Structure as:\n"
            "**Position A (For):** [3+ arguments]\n"
            "**Position B (Against):** [3+ arguments]\n"
            "**Key Evidence:** [for each side]\n"
            "**Balanced Conclusion:** [synthesis]"
            if lang == "en" else
            "قدم تحليلاً متوازناً للنقاش. قسمه إلى:\n"
            "**الموقف أ (مؤيد):** [٣ حجج على الأقل]\n"
            "**الموقف ب (معارض):** [٣ حجج على الأقل]\n"
            "**الأدلة الرئيسية:** [لكل جانب]\n"
            "**خلاصة متوازنة:** [تجميع]"
        ),
        "timeline": (
            "Create a detailed chronological timeline using this format:\n"
            "📅 **[Date/Period]** — **[Event/Development]**\n"
            "   → [Significance and impact]\n"
            "Connect events to show cause and effect relationships."
            if lang == "en" else
            "أنشئ تسلسلاً زمنياً مفصلاً بهذا التنسيق:\n"
            "📅 **[التاريخ/الفترة]** — **[الحدث/التطور]**\n"
            "   → [الأهمية والتأثير]\n"
            "اربط الأحداث لتوضيح علاقات السبب والنتيجة."
        ),
        "chat": "",
    }

    return base + mode_prompts.get(mode, "")

# ==============================================================================
# 💬  CHAT LOGIC
# ==============================================================================
# The core generator that streams LLM tokens to the Gradio UI while persisting to the database.

def chat_logic(message, history, username):
    """
    Manages the core LLM conversation loop using a Generator for token streaming.
    Synchronizes the asynchronous Gradio UI state with the persistent SQLite backend to prevent cross-session memory leaks and prompt contamination.
    """
    """Main chat handler — streaming generator."""
    if not username:
        yield "⚠️ Please log in with Hugging Face to start chatting."
        return

    if not is_valid_username(username):
        yield "⚠️ Invalid username."
        return

    if not rate_limit_ok(username):
        yield "⏳ Rate limit reached — please wait a minute."
        return

    if not message or not message.strip():
        yield "💬 Type a message to start chatting!"
        return

    message = message.strip()
    # If the frontend history is empty, the user cleared the chat.
    # Force a new backend session so the LLM truly forgets previous context.
    if len(history) == 0:
        session_id = create_session(username)
        ctx = get_user_context(username)
        ctx["session_id"] = session_id
        db_history = []
    else:
        session_id = get_or_create_user_session(username)
        db_history = load_session_messages(session_id, limit=MAX_HISTORY_TURNS * 2)

    # Save user message
    save_message(session_id, "user", message)

    # Detect language and mode
    lang = detect_language(message)
    mode = detect_mode(message)

    # Build contextual query
    ctx_query = build_contextual_query(message, db_history)

    # Retrieve documents
    has_docs = False
    context  = ""
    sources  = []
    enabled  = get_enabled_docs(username)

    if enabled:
        search_q = rewrite_query(ctx_query, lang)
        results = hybrid_search(search_q, username, top_k=6)
        if results:
            has_docs = True
            context_parts = []
            seen_sources = set()
            total_chars = 0
            for doc in results:
                chunk_text = doc.page_content
                if total_chars + len(chunk_text) > MAX_CONTEXT_CHARS:
                    break
                context_parts.append(chunk_text)
                total_chars += len(chunk_text)
                src = doc.metadata.get("source", "")
                if src and src not in seen_sources:
                    sources.append(src)
                    seen_sources.add(src)
            context = "\n\n---\n\n".join(context_parts)
    elif requires_document(message) and not enabled:
        note = ("📁 No documents loaded. Upload files for document-based answers.\n\n"
                if lang == "en" else
                "📁 لا توجد مستندات محمّلة. ارفع ملفات للحصول على إجابات من المستندات.\n\n")
        yield note
        return

    # Build messages
    system_prompt = build_system_prompt(mode, lang, has_docs, context)
    session_summary = get_session_summary(session_id)
    if session_summary:
        system_prompt += f"\n\nConversation summary so far: {session_summary}"

    messages = [{"role": "system", "content": system_prompt}]

    # Add history
    for h in db_history:
        role = h.get("role", "")
        content = _extract_text(h.get("content", ""))
        if role in ("user", "assistant") and content:
            messages.append({"role": role, "content": content[:2000]})

    messages.append({"role": "user", "content": message})

    # Source citation
    if sources:
        src_note = "📚 " + " | ".join(f"`{s}`" for s in sources[:4])
    else:
        src_note = ""

    final_reply = ""  # ✅ FIX: Initialize before try block

    yield "⏳"

    try:
        stream = client.chat_completion(
            messages, max_tokens=MAX_RESPONSE_TOKENS, stream=True, temperature=0.7, top_p=0.95,
        )
        partial = ""
        is_arabic = False
        lang_checked = False
        word_count = 0

        for chunk in stream:
            if chunk.choices and chunk.choices[0].delta.content:
                raw = _extract_text(chunk.choices[0].delta.content)
                partial += raw
                word_count += raw.count(" ") + raw.count("\n")

                # Detect language after enough text
                if not lang_checked and len(partial) > 30:
                    is_arabic = (detect_language(partial) == "ar")
                    lang_checked = True

                # Yield every ~3 words for smooth streaming
                if word_count >= 3:
                    word_count = 0
                    display = format_arabic_response(partial) if is_arabic else partial
                    yield display

        # Final yield
        display = format_arabic_response(partial) if is_arabic else partial
        if src_note:
            display = display + "\n\n" + src_note
        yield display
        final_reply = display

    except Exception as e:
        err = str(e)
        log.error(f"Chat completion error: {err}")
        err_msg = ""
        if "410" in err or "deprecated" in err.lower():
            err_msg = "🔴 Model unavailable. Check REPO_ID."
        elif "503" in err:
            err_msg = "🟡 Server loading — please wait 60 s and retry."
        elif "401" in err or "403" in err:
            err_msg = "🔑 Auth error — check HF_TOKEN."
        else:
            err_msg = f"⚠️ {err}"
            
        try:
            display = format_arabic_response(partial) if is_arabic else partial
            final_reply = display + f"\n\n[{err_msg}]"
        except NameError:
            final_reply = err_msg
            
        yield final_reply

    # Save assistant response
    if final_reply:
        save_message(session_id, "assistant", final_reply, sources=sources)
        save_user_history(username, history + [
            {"role": "user", "content": message},
            {"role": "assistant", "content": final_reply},
        ])
        update_summary_if_needed(session_id, username)

# ==============================================================================
# 🔊  TTS + AUDIO
# ==============================================================================
# Converts text to MP3 using a resilient, multi-tiered fallback architecture (Azure -> HF -> gTTS).

VOICE_MAP = {
    "ar": {"azure": "ar-EG-SalmaNeural", "gtts": "ar"},
    "en": {"azure": "en-US-JennyNeural", "gtts": "en"},
    "fr": {"azure": "fr-FR-DeniseNeural", "gtts": "fr"},
}


def resolve_lang(text: str) -> str:
    """Resolve language code for TTS."""
    lang = detect_language(text)
    return lang if lang in VOICE_MAP else "en"


def tts_to_audio(text: str, lang: str = None) -> tuple[Optional[str], Optional[str]]:
    """
    Converts text to an MP3 audio stream using a multi-tiered fallback architecture.
    Attempts Azure Neural TTS first with SSML tags, gracefully falling back to Hugging Face models or gTTS if the primary API fails or rate-limits.
    """
    """Convert text to audio using Azure TTS (preferred) or gTTS fallback. Returns (filepath, error_msg)."""
    if not text or not text.strip():
        return None, "Empty text provided."

    if lang is None:
        lang = resolve_lang(text)

    # Clean text for TTS
    clean = re.sub(r"[#*`_~\[\](){}<>|]", "", text)
    clean = re.sub(r"https?://\S+", "", clean)
    clean = re.sub(r"\n{2,}", ". ", clean)
    clean = re.sub(r"\n", " ", clean)
    clean = clean.strip()

    if not clean:
        return None, "Text was empty after cleaning."

    # Try Azure TTS first
    if AZURE_TTS_AVAILABLE and AZURE_SPEECH_KEY:
        try:
            voice = VOICE_MAP.get(lang, VOICE_MAP["en"])["azure"]
            config = speechsdk.SpeechConfig(
                subscription=AZURE_SPEECH_KEY,
                region=AZURE_SPEECH_REGION,
            )
            config.speech_synthesis_voice_name = voice
            config.set_speech_synthesis_output_format(
                speechsdk.SpeechSynthesisOutputFormat.Audio16Khz32KBitRateMonoMp3
            )
            tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".mp3", delete=False)
            audio_config = speechsdk.audio.AudioOutputConfig(filename=tmp.name)
            synth = speechsdk.SpeechSynthesizer(
                speech_config=config, audio_config=audio_config
            )
            
            # Wrap English words in SSML for proper pronunciation
            if lang == "ar":
                import xml.sax.saxutils as saxutils
                parts = re.split(r'([a-zA-Z][a-zA-Z0-9\s\-\.]*)', clean)
                ssml_clean = ""
                for i, part in enumerate(parts):
                    if not part: continue
                    escaped_part = saxutils.escape(part)
                    if i % 2 == 1:
                        ssml_clean += f'<lang xml:lang="en-US">{escaped_part}</lang>'
                    else:
                        ssml_clean += escaped_part
                        
                ssml = f"""<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xml:lang="ar-EG">
    <voice name="{voice}">
        <prosody volume="+100%">{ssml_clean}</prosody>
    </voice>
</speak>"""
                result = synth.speak_ssml_async(ssml).get()
            else:
                import xml.sax.saxutils as saxutils
                ssml = f"""<speak version="1.0" xmlns="http://www.w3.org/2001/10/synthesis" xml:lang="en-US">
    <voice name="{voice}">
        <prosody volume="+100%">{saxutils.escape(clean)}</prosody>
    </voice>
</speak>"""
                result = synth.speak_ssml_async(ssml).get()
                
            if result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
                return tmp.name, None
            else:
                log.warning(f"Azure TTS failed ({result.reason}), falling back to alternative TTS engines.")
                pass
        except Exception as e:
            log.error(f"Azure TTS error: {e}")

    # Try HF TTS tier (Only for Arabic, HF English TTS is robotic)
    if lang == "ar":
        try:
            hf_model = TTS_MODEL_AR
            audio_bytes = client.text_to_speech(clean, model=hf_model)
            if audio_bytes:
                tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".mp3", delete=False)
                tmp.write(audio_bytes)
                tmp.close()
                return tmp.name, None
        except Exception as e:
            log.error(f"HF TTS error: {e}")

    # Fallback to gTTS
    if GTTS_AVAILABLE:
        try:
            gtts_lang = VOICE_MAP.get(lang, VOICE_MAP["en"])["gtts"]
            tts = gTTS(text=clean, lang=gtts_lang, slow=False)
            tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".mp3", delete=False)
            tts.save(tmp.name)
            return tmp.name, None
        except Exception as e:
            log.error(f"gTTS error: {e}")
            return None, f"gTTS failed: {str(e)} (HF Spaces might be IP blocked by Google)"
    
    return None, "No TTS engine available (set AZURE_SPEECH_KEY or use gTTS)"

# ==============================================================================
# 🎬  VIDEO GENERATION
# ==============================================================================
# Generates MP4 videos programmatically by drawing slides with Pillow and compositing with MoviePy.

def _font(size=28, bold=False, is_arabic=False):
    """Get a font for text rendering."""
    try:
        from PIL import ImageFont as IF
        import urllib.request
        
        if is_arabic:
            reg_path = "Tajawal-Regular.ttf"
            bold_path = "Tajawal-Bold.ttf"
            
            if not os.path.exists(reg_path):
                try: urllib.request.urlretrieve("https://github.com/google/fonts/raw/main/ofl/tajawal/Tajawal-Regular.ttf", reg_path)
                except: pass
            if not os.path.exists(bold_path):
                try: urllib.request.urlretrieve("https://github.com/google/fonts/raw/main/ofl/tajawal/Tajawal-Bold.ttf", bold_path)
                except: pass

            font_paths = [
                bold_path if bold else reg_path,
                "/usr/share/fonts/truetype/noto/NotoSansArabic-Bold.ttf" if bold else "/usr/share/fonts/truetype/noto/NotoSansArabic-Regular.ttf",
                "/usr/share/fonts/google-noto/NotoSansArabic-Bold.ttf" if bold else "/usr/share/fonts/google-noto/NotoSansArabic-Regular.ttf",
            ]
        else:
            font_paths = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/dejavu-sans-fonts/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/dejavu-sans-fonts/DejaVuSans.ttf",
            ]
            
        for fp in font_paths:
            if os.path.exists(fp):
                return IF.truetype(fp, size)
        return IF.load_default()
    except Exception as e:
        log.error(f"Font error: {e}")
        return None


def _clean_markdown(text):
    """Strip markdown formatting for clean display."""
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)   # **bold**
    text = re.sub(r'\*([^*]+)\*', r'\1', text)       # *italic*
    text = re.sub(r'__([^_]+)__', r'\1', text)       # __bold__
    text = re.sub(r'_([^_]+)_', r'\1', text)         # _italic_
    text = re.sub(r'#{1,6}\s*', '', text)             # headers
    text = re.sub(r'`([^`]+)`', r'\1', text)          # inline code
    text = re.sub(r'```[\s\S]*?```', '', text)        # code blocks
    text = re.sub(r'[\-\*]\s+', '• ', text)           # bullet points
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)  # links
    return text.strip()


_RESHAPER_CACHE = None


def chunk_text(text, max_chars=350):
    """Split text into chunks for slide display, respecting numbered sections."""
    # First try to split on numbered sections (1. or 1-)
    section_splits = re.split(r'(?=\n\s*\d+[\.\-\)]\s)', text)
    
    chunks = []
    for section in section_splits:
        section = section.strip()
        if not section:
            continue
        if len(section) <= max_chars:
            chunks.append(section)
        else:
            # Split long sections by sentences
            sentences = re.split(r'(?<=[.!?؟])\s+', section)
            current = ""
            for s in sentences:
                if len(current) + len(s) > max_chars and current:
                    chunks.append(current.strip())
                    current = s
                else:
                    current = current + " " + s if current else s
            if current.strip():
                chunks.append(current.strip())
    
    # Fallback: if we got 0 or 1 chunk, force-split by paragraphs
    if len(chunks) <= 1 and len(text) > max_chars:
        chunks = []
        paragraphs = re.split(r'\n', text)
        current = ""
        for p in paragraphs:
            p = p.strip()
            if not p:
                continue
            if len(current) + len(p) + 1 > max_chars and current:
                chunks.append(current.strip())
                current = p
            else:
                current = (current + " " + p).strip() if current else p
        if current.strip():
            chunks.append(current.strip())
    
    return chunks if chunks else [text[:max_chars]]


# ── Color constants ──────────────────────────────────────────────────────────
SLIDE_COLORS = [
    {"bg1": (15, 12, 30),  "bg2": (30, 20, 55),  "accent": (255, 107, 61),  "text": (236, 238, 245)},
    {"bg1": (10, 18, 30),  "bg2": (20, 35, 55),  "accent": (108, 138, 255), "text": (236, 238, 245)},
    {"bg1": (12, 25, 18),  "bg2": (22, 45, 32),  "accent": (34, 212, 122),  "text": (236, 238, 245)},
    {"bg1": (25, 15, 10),  "bg2": (45, 28, 18),  "accent": (240, 195, 70),  "text": (236, 238, 245)},
    {"bg1": (20, 10, 25),  "bg2": (38, 18, 48),  "accent": (200, 120, 255), "text": (236, 238, 245)},
]


def _draw_gradient_bg(draw, width, height, color1, color2):
    """Draw a vertical gradient background."""
    for y in range(height):
        ratio = y / height
        r = int(color1[0] * (1 - ratio) + color2[0] * ratio)
        g = int(color1[1] * (1 - ratio) + color2[1] * ratio)
        b = int(color1[2] * (1 - ratio) + color2[2] * ratio)
        draw.line([(0, y), (width, y)], fill=(r, g, b))


def _draw_rounded_rect(draw, xy, radius, fill):
    """Draw a rounded rectangle."""
    x1, y1, x2, y2 = xy
    draw.rectangle([x1 + radius, y1, x2 - radius, y2], fill=fill)
    draw.rectangle([x1, y1 + radius, x2, y2 - radius], fill=fill)
    draw.pieslice([x1, y1, x1 + 2 * radius, y1 + 2 * radius], 180, 270, fill=fill)
    draw.pieslice([x2 - 2 * radius, y1, x2, y1 + 2 * radius], 270, 360, fill=fill)
    draw.pieslice([x1, y2 - 2 * radius, x1 + 2 * radius, y2], 90, 180, fill=fill)
    draw.pieslice([x2 - 2 * radius, y2 - 2 * radius, x2, y2], 0, 90, fill=fill)


def _wrap_text_to_lines(text, font, max_width, draw, is_rtl=False):
    """Wrap text to fit within max_width."""
    words = text.split()
    lines = []
    current_line = ""
    direction = "rtl" if is_rtl else "ltr"
    for word in words:
        test = current_line + " " + word if current_line else word
        try:
            bbox = draw.textbbox((0, 0), test, font=font, direction=direction)
            w = bbox[2] - bbox[0]
        except Exception:
            w = len(test) * 10
        if w <= max_width:
            current_line = test
        else:
            if current_line:
                lines.append(current_line)
            current_line = word
    if current_line:
        lines.append(current_line)
    return lines


def create_slide_image(text, slide_num, total_slides, lang="en", width=1280, height=720):
    """Create a modern slide image with text and decorative elements."""
    if not PIL_AVAILABLE:
        return None

    try:
        # Clean markdown from text
        text = _clean_markdown(text)
        is_rtl = lang == "ar"

        colors = SLIDE_COLORS[slide_num % len(SLIDE_COLORS)]
        
        # Try AI background first
        ai_bg = generate_slide_bg(text, width=width, height=height)
        if ai_bg:
            img = ai_bg
            draw = ImageDraw.Draw(img)
        else:
            img = Image.new("RGB", (width, height))
            draw = ImageDraw.Draw(img)
            _draw_gradient_bg(draw, width, height, colors["bg1"], colors["bg2"])

        accent = colors["accent"]
        accent_dim = tuple(max(0, c - 80) for c in accent)

        # Top accent bar (gradient feel)
        for x in range(width):
            ratio = x / width
            r = int(accent[0] * ratio + accent_dim[0] * (1 - ratio))
            g = int(accent[1] * ratio + accent_dim[1] * (1 - ratio))
            b = int(accent[2] * ratio + accent_dim[2] * (1 - ratio))
            draw.line([(x, 0), (x, 4)], fill=(r, g, b))

        # Decorative corner circles (subtle)
        for cx, cy, radius in [(width - 80, 80, 60), (100, height - 100, 40)]:
            for r in range(radius, 0, -1):
                alpha = int(15 * (r / radius))
                c = tuple(min(255, ch + alpha) for ch in colors["bg2"])
                draw.ellipse([cx - r, cy - r, cx + r, cy + r], outline=c)

        # Progress bar at bottom
        bar_y = height - 8
        draw.rectangle([0, bar_y, width, height], fill=(30, 30, 40))
        progress_width = int(width * (slide_num + 1) / total_slides)
        draw.rectangle([0, bar_y, progress_width, height], fill=accent)

        # Slide number badge
        small_font = _font(16, is_arabic=is_rtl)
        if small_font:
            num_text = f"{slide_num + 1} / {total_slides}"
            badge_w = 90
            badge_x = width - badge_w - 30
            _draw_rounded_rect(draw, (badge_x, 16, badge_x + badge_w, 42), 6, accent)
            draw.text((badge_x + 12, 20), num_text, fill=(220, 220, 230), font=small_font)

        # Text layout
        title_font = _font(30, bold=True, is_arabic=is_rtl)
        body_font  = _font(22, is_arabic=is_rtl)

        lines_raw = text.split("\n")
        if len(lines_raw) == 1:
            title = ""
            body = text
        else:
            title = lines_raw[0]
            body  = "\n".join(lines_raw[1:])

        margin = 70
        x_pos = width - margin if is_rtl else margin
        anchor = "ra" if is_rtl else "la"
        direction = "rtl" if is_rtl else "ltr"
        max_text_w = width - margin * 2

        # Draw title with underline accent
        y_offset = 60
        if title and title_font:
            title_lines = _wrap_text_to_lines(title, title_font, max_text_w, draw, is_rtl=is_rtl)
            for tl in title_lines[:2]:
                draw.text((x_pos, y_offset), tl, fill=accent, font=title_font, anchor=anchor, direction=direction)
                y_offset += 42
            # Underline below title
            line_x1 = margin
            line_x2 = min(margin + 300, width - margin)
            if is_rtl:
                line_x1 = max(width - margin - 300, margin)
                line_x2 = width - margin
            draw.rectangle([line_x1, y_offset + 4, line_x2, y_offset + 6], fill=accent)
            y_offset += 20

        # Draw body
        if body and body_font:
            body_lines = _wrap_text_to_lines(body, body_font, max_text_w, draw, is_rtl=is_rtl)
            for bl in body_lines:
                if y_offset > height - 50:
                    break
                draw.text((x_pos, y_offset), bl, fill=colors["text"], font=body_font, anchor=anchor, direction=direction)
                y_offset += 34

        # Watermark
        wm_font = _font(14)
        if wm_font:
            draw.text((margin, height - 28), "Phoenix Eduplan v4.1", fill=(80, 80, 100), font=wm_font)

        tmp = tempfile.NamedTemporaryFile(dir=".", suffix=".png", delete=False)
        img.save(tmp.name, "PNG")
        return tmp.name
    except Exception as e:
        log.error(f"Slide creation error: {e}")
        return None


def text_to_video(text, lang=None):
    """
    Programmatically orchestrates the rendering of a video file (.mp4) from a textual prompt.
    Chunks the script into visual slides, draws the UI components via PIL, synthesizes audio for each chunk, and composites the final Directed Acyclic Graph (DAG) using MoviePy.
    """
    """Generate a video from text with slides and audio narration."""
    if not MOVIEPY_AVAILABLE:
        return None, "❌ moviepy not available."
    if not PIL_AVAILABLE:
        return None, "❌ Pillow not available."

    if lang is None:
        lang = resolve_lang(text)

    # Clean markdown before creating slides
    sanitized_text = _clean_markdown(text)
    slides_text = chunk_text(sanitized_text, max_chars=350)
    if not slides_text:
        return None, "❌ No content for video."

    clips = []
    temp_files = []

    try:
        for i, slide_text in enumerate(slides_text):
            # Create slide image
            img_path = create_slide_image(slide_text, i, len(slides_text), lang=lang)
            if not img_path:
                continue
            temp_files.append(img_path)

            # Create audio for slide
            audio_path, _ = tts_to_audio(slide_text, lang)

            if audio_path:
                temp_files.append(audio_path)
                audio_clip = AudioFileClip(audio_path)
                duration = audio_clip.duration + 0.5
                img_clip = ImageClip(img_path).set_duration(duration)
                img_clip = img_clip.fadein(0.5).fadeout(0.5)
                img_clip = img_clip.set_audio(audio_clip)
            else:
                word_count = len(slide_text.split())
                duration = max(3, word_count / 2.5)
                img_clip = ImageClip(img_path).set_duration(duration)
                img_clip = img_clip.fadein(0.5).fadeout(0.5)

            clips.append(img_clip)

        if not clips:
            return None, "❌ Failed to create slides."

        has_audio = any(c.audio is not None for c in clips)
        video = concatenate_videoclips(clips, method="compose")
        tmp_video = tempfile.NamedTemporaryFile(dir=".", suffix=".mp4", delete=False)
        video.write_videofile(
            tmp_video.name, fps=24, codec="libx264",
            audio_codec="aac" if has_audio else None,
            audio=has_audio,
            preset="ultrafast",
            threads=4,
            logger=None,
        )
        duration = video.duration
        video.close()

        return tmp_video.name, f"✅ Video created — {len(clips)} slides, {duration:.0f}s"
    except Exception as e:
        log.error(f"Video creation error: {e}")
        return None, f"❌ Video error: {e}"
    finally:
        for c in clips:
            try:
                c.close()
            except:
                pass
        for tf in temp_files:
            try:
                if os.path.exists(tf):
                    os.remove(tf)
            except:
                pass


def build_script_from_docs(username, lang="en"):
    """Build a comprehensive presentation script from ALL user document pages."""
    enabled = get_enabled_docs(username)
    if not enabled:
        return ""

    ctx = get_user_context(username)
    doc_store = ctx["doc_store"]

    # Read ALL chunks from ALL active documents — do not skip any
    parts = []
    for did in ctx["active_docs"]:
        chunks = doc_store.get(did, [])
        for chunk in chunks:  # ALL chunks, not just first 10
            parts.append(chunk.page_content)

    full_text = "\n\n".join(parts)
    # Allow up to 15000 chars for comprehensive coverage
    if len(full_text) > 15000:
        full_text = full_text[:15000]
        full_text = full_text.rsplit(' ', 1)[0]
        full_text += "\n[Content truncated]"

    if lang == "auto":
        lang = "ar" if detect_language(full_text) == "ar" else "en"

    try:
        if lang == "ar":
            sys_msg = (
                "أنت معلم مصري خبير وتشرح المادة لطلابك بطريقة تفاعلية وممتعة جداً باللهجة المصرية. "
                "يجب أن تغطي كل النقاط والمواضيع الموجودة في المحتوى دون إسقاط أي شيء."
            )
            prompt = (
                "اقرأ كل المحتوى التالي بعناية شديدة صفحة بصفحة، ثم أنشئ نص عرض تقديمي شامل يغطي كل نقطة. "
                "تحدث وكأنك معلم مصري يشرح لطلابه بحماس. استخدم جملاً انتقالية سلسة بين المواضيع وتحدث باللهجة المصرية المفهومة والمبسطة. "
                "لا تترك أي معلومة أو فقرة أو نقطة دون تغطيتها! "
                "تعليمات صارمة: أبقِ كل المصطلحات التقنية بالإنجليزية كما هي "
                "(مثال صحيح: 'تقنية الـ Virtualization' — مثال خاطئ: 'تقنية الافتراضية'). "
                "لا تستخدم أي تنسيق Markdown مثل ** أو ## أو * في النص. "
                "لا تستخدم قوائم مرقمة (1, 2, 3) أبداً! اكتب النص كقصة أو محاضرة متصلة وممتعة:\n\n" + full_text
            )
        else:
            sys_msg = (
                "You are an expert teacher explaining the material to your students in an engaging way. "
                "You must cover EVERY topic and point from the provided content without skipping anything."
            )
            prompt = (
                "Read ALL the following content carefully page by page, then create a comprehensive "
                "presentation script that covers EVERY point, topic, and detail mentioned. "
                "Speak like a teacher giving an engaging lecture. Use smooth transitions between topics. "
                "Do NOT skip any information! Do NOT use any Markdown formatting like ** or ## or * — write plain text only. "
                "Do NOT use numbered lists (1, 2, 3)! Write the text as a continuous, engaging lecture:\n\n" + full_text
            )
            
        messages = [
            {"role": "system", "content": sys_msg},
            {"role": "user", "content": prompt}
        ]
        resp = client.chat_completion(messages, max_tokens=3000, temperature=0.4)
        script = _extract_text(resp.choices[0].message.content)
        # Clean any leftover markdown
        script = _clean_markdown(script)
        return script.strip() if script else full_text[:5000]
    except Exception as e:
        log.error(f"Script generation error: {e}")
        return full_text[:5000]


# ==============================================================================
# 🎨  AI IMAGE GENERATION (FLUX)
# ==============================================================================
# Uses Hugging Face Inference endpoints to generate cinematic backgrounds for video slides.

USE_FLUX_BGS = False  # Set to True to enable FLUX AI backgrounds for video slides

def generate_slide_bg(topic_text, width=1280, height=720):
    """Generate an AI background image for a video slide using FLUX."""
    if not PIL_AVAILABLE or not USE_FLUX_BGS:
        return None
    try:
        # Create a cinematic prompt from the topic
        prompt = f"Dark futuristic digital illustration of {topic_text[:80]}, abstract tech background, cinematic lighting, deep blue and orange tones, 16:9 aspect ratio, no text"
        img_client = InferenceClient(token=HF_TOKEN)
        result = img_client.text_to_image(
            prompt,
            model=IMAGE_MODEL,
            width=width,
            height=height,
        )
        if result:
            # Apply dark overlay for text readability
            overlay = Image.new("RGBA", result.size, (10, 8, 20, 160))
            result = result.convert("RGBA")
            result = Image.alpha_composite(result, overlay)
            return result.convert("RGB")
    except Exception as e:
        log.warning(f"AI image generation failed (using fallback): {e}")
    return None


# ==============================================================================
# 📋  STUDY PLAN GENERATOR
# ==============================================================================
# Injects RAG chunks into a bounded prompt to output a structured markdown study schedule.

def generate_study_plan(username, duration="2 weeks", lang_choice="auto"):
    """
    Acts as an autonomous study planner by injecting all of the student's FAISS chunks directly into the LLM context window.
    Dynamically bounds the LLM generation with a strict system prompt to force output into a structured, day-by-day markdown schedule.
    """
    """Generate a personalized study plan from uploaded documents."""
    enabled = get_enabled_docs(username)
    if not enabled:
        return "⚠️ Upload documents first to generate a study plan."

    ctx = get_user_context(username)
    doc_store = ctx["doc_store"]

    parts = []
    for did in ctx["active_docs"]:
        chunks = doc_store.get(did, [])
        for chunk in chunks:
            parts.append(chunk.page_content)

    full_text = "\n\n".join(parts)
    if len(full_text) > 12000:
        full_text = full_text[:12000].rsplit(' ', 1)[0]

    lang = "ar" if lang_choice == "ar" else ("ar" if detect_language(full_text) == "ar" and lang_choice == "auto" else "en")

    try:
        if lang == "ar":
            sys_msg = (
                "أنت مخطط دراسي ذكي. أنشئ خطة دراسة مفصلة وعملية "
                "باللهجة المصرية الممتعة والواضحة. ابقِ المصطلحات التقنية بالإنجليزية."
            )
            prompt = (
                f"اقرأ المحتوى التالي بعناية وأنشئ خطة دراسة مفصلة لمدة {duration}. "
                "رد باللهجة المصرية بأسلوب مشجع وبناء. "
                "أضف جلسات مراجعة ونقاط فحص ذاتي. "
                f"استخدم الرموز 📚🎯✅🔁 لتنظيم الخطة.\n\nالمحتوى:\n{full_text}"
            )
        else:
            sys_msg = (
                "You are an expert study planner. Create a detailed, actionable study plan "
                "with daily goals, review sessions, and self-test checkpoints."
            )
            prompt = (
                f"Read the following content carefully and create a detailed study plan for {duration}. "
                "Structure it day-by-day with clear daily objectives. "
                "Include review sessions every few days and self-test checkpoints. "
                "Rate topic difficulty (🟢 Easy / 🟡 Medium / 🔴 Hard). "
                f"Use emojis 📚🎯✅🔁 for organization.\n\nContent:\n{full_text}"
            )

        messages = [
            {"role": "system", "content": sys_msg},
            {"role": "user", "content": prompt}
        ]
        resp = client.chat_completion(messages, max_tokens=3000, temperature=0.5)
        plan = _extract_text(resp.choices[0].message.content)
        return plan.strip() if plan else "❌ Failed to generate study plan."
    except Exception as e:
        log.error(f"Study plan generation error: {e}")
        return f"❌ Error: {e}"


# ==============================================================================
# ✍️  ESSAY GRADER
# ==============================================================================
# Uses prompt engineering to force the LLM to output deterministically formatted grading rubrics.

def grade_essay(essay_text, rubric="", lang_choice="auto"):
    """
    Executes a structured evaluation of student writing using Llama-3.
    Uses prompt-engineering techniques to force the LLM to output a deterministically formatted grading rubric, highlighting strengths and specific areas for improvement.
    """
    """Grade an essay using AI with detailed feedback."""
    if not essay_text or not essay_text.strip():
        return "⚠️ Please paste your essay text."

    if len(essay_text) > 15000:
        return "⚠️ Essay is too long. Please limit to 15,000 characters."

    lang = "ar" if lang_choice == "ar" else ("ar" if detect_language(essay_text) == "ar" and lang_choice == "auto" else "en")

    try:
        if lang == "ar":
            sys_msg = (
                "أنت معلم مصري خبير في تقييم المقالات. قيم المقال بدقة وموضوعية "
                "مع تقديم ملاحظات بنّاءة ومفصلة باللهجة المصرية. ابقِ المصطلحات التقنية بالإنجليزية."
            )
            rubric_note = f"\n\nمعايير التقييم: {rubric}" if rubric else ""
            prompt = (
                f"قيم هذا المقال وقدم:\n"
                "🎯 **الدرجة الإجمالية:** (A/B/C/D/F)\n"
                "📊 **تفصيل الدرجات:**\n"
                "   • المحتوى: _/10\n   • التنظيم: _/10\n   • اللغة: _/10\n   • التفكير النقدي: _/10\n"
                "✅ **نقاط القوة:**\n"
                "⚠️ **نقاط التحسين:**\n"
                "💡 **اقتراحات محددة:**\n"
                f"📝 **إعادة كتابة مقترحة لأضعف فقرة:**{rubric_note}\n\nالمقال:\n{essay_text}"
            )
        else:
            sys_msg = (
                "You are an expert essay grader. Evaluate the essay accurately and objectively "
                "with detailed, constructive feedback."
            )
            rubric_note = f"\n\nGrading Rubric: {rubric}" if rubric else ""
            prompt = (
                f"Grade this essay and provide:\n"
                "🎯 **Overall Grade:** (A/B/C/D/F)\n"
                "📊 **Score Breakdown:**\n"
                "   • Content & Ideas: _/10\n   • Organization: _/10\n   • Language & Style: _/10\n   • Critical Thinking: _/10\n"
                "✅ **Strengths:**\n"
                "⚠️ **Areas for Improvement:**\n"
                "💡 **Specific Suggestions:**\n"
                f"📝 **Suggested Rewrite of Weakest Paragraph:**{rubric_note}\n\nEssay:\n{essay_text}"
            )

        messages = [
            {"role": "system", "content": sys_msg},
            {"role": "user", "content": prompt}
        ]
        resp = client.chat_completion(messages, max_tokens=2500, temperature=0.4)
        result = _extract_text(resp.choices[0].message.content)
        return result.strip() if result else "❌ Failed to grade essay."
    except Exception as e:
        log.error(f"Essay grading error: {e}")
        return f"❌ Error: {e}"


# ==============================================================================
# 📚  VOCABULARY EXTRACTOR
# ==============================================================================
# Extracts key terms and definitions from academic texts for flashcard generation.

def extract_vocabulary(username, lang_choice="auto"):
    """Extract key terms and definitions from uploaded documents."""
    enabled = get_enabled_docs(username)
    if not enabled:
        return "⚠️ Upload documents first."

    ctx = get_user_context(username)
    doc_store = ctx["doc_store"]

    parts = []
    for did in ctx["active_docs"]:
        chunks = doc_store.get(did, [])
        for chunk in chunks:
            parts.append(chunk.page_content)

    full_text = "\n\n".join(parts)
    if len(full_text) > 10000:
        full_text = full_text[:10000].rsplit(' ', 1)[0]

    lang = "ar" if lang_choice == "ar" else ("ar" if detect_language(full_text) == "ar" and lang_choice == "auto" else "en")

    try:
        if lang == "ar":
            sys_msg = "أنت خبير في استخراج المصطلحات الرئيسية والمفاهيم من النصوص الأكاديمية وتشرحها باللهجة المصرية المبسطة. ابقِ المصطلحات التقنية بالإنجليزية."
            prompt = (
                "استخرج جميع المصطلحات الرئيسية والمفاهيم المهمة من النص. لكل مصطلح قدم:\n"
                "📌 **المصطلح:** [Term]\n"
                "📖 **التعريف:** [Definition in English or clear Egyptian Arabic]\n"
                "💡 **مثال:** [Example]\n"
                "🔗 **مفاهيم مرتبطة:** [Related]\n---\n"
                f"\nالنص:\n{full_text}"
            )
        else:
            sys_msg = "You are an expert at extracting key terms and concepts from academic texts."
            prompt = (
                "Extract ALL key terms and important concepts from this text. For each term provide:\n"
                "📌 **Term:** [name]\n"
                "📖 **Definition:** [clear definition]\n"
                "💡 **Example:** [practical example]\n"
                "🔗 **Related Concepts:** [connections]\n---\n"
                f"\nText:\n{full_text}"
            )

        messages = [
            {"role": "system", "content": sys_msg},
            {"role": "user", "content": prompt}
        ]
        resp = client.chat_completion(messages, max_tokens=3000, temperature=0.3)
        result = _extract_text(resp.choices[0].message.content)
        return result.strip() if result else "❌ Failed to extract vocabulary."
    except Exception as e:
        log.error(f"Vocabulary extraction error: {e}")
        return f"❌ Error: {e}"


# ==============================================================================
# 📊  ANALYTICS
# ==============================================================================
# Aggregates user database metrics to display a personalized learning dashboard.

def get_analytics(username):
    """Generate analytics for the user's learning activity."""
    if not is_valid_username(username):
        return "⚠️ Enter your name first."

    try:
        with get_db() as conn:
            # Total sessions
            sessions = conn.execute(
                "SELECT COUNT(*) FROM sessions WHERE username=?",
                (sanitize_username(username),)
            ).fetchone()[0]

            # Total messages
            total_msgs = conn.execute(
                "SELECT COUNT(*) FROM messages m JOIN sessions s ON m.session_id = s.id WHERE s.username=?",
                (sanitize_username(username),)
            ).fetchone()[0]

            # User messages vs assistant messages
            user_msgs = conn.execute(
                "SELECT COUNT(*) FROM messages m JOIN sessions s ON m.session_id = s.id WHERE s.username=? AND m.role='user'",
                (sanitize_username(username),)
            ).fetchone()[0]

            assistant_msgs = total_msgs - user_msgs

            # Recent messages for topic analysis
            recent = conn.execute(
                "SELECT m.content FROM messages m JOIN sessions s ON m.session_id = s.id "
                "WHERE s.username=? AND m.role='user' ORDER BY m.id DESC LIMIT 20",
                (sanitize_username(username),)
            ).fetchall()

        # Document stats
        ctx = get_user_context(username)
        total_docs = len(ctx["doc_store"])
        active_docs = len(ctx["active_docs"])
        total_chunks = sum(len(v) for v in ctx["doc_store"].values())

        # Build analytics report
        report = f"""## 📊 Your Learning Analytics

### 📈 Activity Overview
| Metric | Value |
|--------|-------|
| 💬 Total Sessions | {sessions} |
| 📝 Your Messages | {user_msgs} |
| 🤖 AI Responses | {assistant_msgs} |
| 📁 Documents Loaded | {total_docs} |
| ✅ Active Documents | {active_docs} |
| 📦 Total Text Chunks | {total_chunks} |

### 🎯 Recent Topics
"""
        if recent:
            topics = [r[0][:80] + "..." if len(r[0]) > 80 else r[0] for r in recent[:10]]
            for i, topic in enumerate(topics, 1):
                report += f"{i}. {topic}\n"
        else:
            report += "_No recent activity yet._\n"

        # Engagement score
        engagement = min(100, int((total_msgs / max(sessions, 1)) * 10))
        bar_filled = "█" * (engagement // 5)
        bar_empty = "░" * (20 - engagement // 5)
        report += f"\n### 🔥 Engagement Score\n`{bar_filled}{bar_empty}` **{engagement}%**\n"

        if engagement < 30:
            report += "\n> 💡 _Tip: Try uploading more documents and asking deeper questions!_"
        elif engagement < 70:
            report += "\n> 💡 _Tip: Try using study plans and practice problems to level up!_"
        else:
            report += "\n> 🌟 _You're a power learner! Keep up the great work!_"

        return report
    except Exception as e:
        log.error(f"Analytics error: {e}")
        return f"❌ Error generating analytics: {e}"


# ==============================================================================
# 📝  AUTO DOCUMENT SUMMARY
# ==============================================================================
# Creates a concise 3-sentence summary of any uploaded document automatically.

def auto_summarize_document(doc_text, lang="en"):
    """Generate an automatic summary when a document is uploaded."""
    if not doc_text or len(doc_text.strip()) < 100:
        return ""
    try:
        text = doc_text[:5000]
        if lang == "ar":
            prompt = f"لخص هذا المحتوى في ٣-٤ جمل قصيرة باللهجة المصرية. ابقِ المصطلحات بالإنجليزية:\n{text}"
        else:
            prompt = f"Summarize this content in 3-4 concise sentences:\n{text}"
        result = client.text_generation(prompt, max_new_tokens=200, temperature=0.3)
        return result.strip() if result else ""
    except Exception as e:
        log.warning(f"Auto-summary failed: {e}")
        return ""


def generate_audio_only(text, username, lang_choice="auto"):
    """Generate audio from text or documents."""
    if not text or not text.strip():
        docs_text = build_script_from_docs(username, lang=lang_choice)
        if docs_text:
            text = docs_text
            lang = resolve_lang(text) if lang_choice == "auto" else lang_choice
        else:
            return None, "⚠️ Enter text or upload documents first."
    else:
        lang = resolve_lang(text) if lang_choice == "auto" else lang_choice

    audio_path, err = tts_to_audio(text, lang)
    if audio_path:
        return audio_path, "✅ Audio generated successfully."
    return None, f"❌ {err}"


def generate_video_only(text, username, lang_choice="auto"):
    """Generate video from text or documents."""
    if not text or not text.strip():
        docs_text = build_script_from_docs(username, lang=lang_choice)
        if docs_text:
            text = docs_text
            lang = resolve_lang(text) if lang_choice == "auto" else lang_choice
        else:
            return None, "⚠️ Enter text or upload documents first."
    else:
        lang = resolve_lang(text) if lang_choice == "auto" else lang_choice

    video_path, msg = text_to_video(text, lang)
    return video_path, msg

# ==============================================================================
# 🎨  CSS
# ==============================================================================
# The master stylesheet defining the dark mode glassmorphism frontend for the Gradio UI.

DARK_CSS = r"""
@import url('https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;800&family=DM+Sans:ital,wght@0,300;0,400;0,500;0,600;1,400&family=Noto+Naskh+Arabic:wght@400;500;700&display=swap');

:root {
    --accent: #ff6b3d;
    --accent-glow: rgba(255, 107, 61, 0.35);
    --gold: #f0c346;
    --flame1: #ff4500;
    --flame2: #ff6b3d;
    --flame3: #ffa500;
    --flame4: #ffd700;
    --bg-deep: #0a0a0f;
    --bg-card: rgba(18, 14, 22, 0.85);
    --bg-surface: rgba(28, 22, 35, 0.7);
    --text-primary: #f0ece4;
    --text-muted: #9e8f84;
    --border-subtle: rgba(255, 107, 61, 0.12);
    --transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
}

/* ── Global overrides ─────────────────────────────────── */
.gradio-container {
    background: var(--bg-deep) !important;
    font-family: 'DM Sans', sans-serif !important;
    max-width: 100% !important;
}

.login-container {
    max-width: 480px !important;
    margin: 12vh auto 5vh auto !important;
    padding: 3rem 2rem !important;
    background: rgba(28, 22, 35, 0.4) !important;
    border: 1px solid rgba(255, 107, 61, 0.2) !important;
    border-radius: 24px !important;
    box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.5), 0 0 40px rgba(255, 107, 61, 0.1) !important;
    backdrop-filter: blur(20px) saturate(180%);
    -webkit-backdrop-filter: blur(20px) saturate(180%);
    text-align: center;
    position: relative;
    overflow: hidden;
}

.login-container::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0; height: 4px;
    background: linear-gradient(90deg, var(--gold), var(--accent), var(--flame1));
}

.login-header-text {
    font-family: 'Outfit', sans-serif !important;
    font-size: 2.5rem !important;
    font-weight: 800 !important;
    background: linear-gradient(135deg, #fff, var(--gold)) !important;
    -webkit-background-clip: text !important;
    -webkit-text-fill-color: transparent !important;
    margin-bottom: 0.5rem !important;
    letter-spacing: -0.5px;
}

.login-subtitle {
    color: var(--text-muted) !important;
    font-size: 1.1rem !important;
    margin-bottom: 2.5rem !important;
}

.custom-input input {
    background: rgba(28, 22, 35, 0.6) !important;
    border: 1px solid rgba(255, 107, 61, 0.3) !important;
    color: var(--text-primary) !important;
    border-radius: 12px !important;
    padding: 0.8rem 1rem !important;
    transition: var(--transition) !important;
    font-family: 'DM Sans', sans-serif !important;
}

.custom-input input:focus {
    border-color: var(--accent) !important;
    box-shadow: 0 0 15px rgba(255, 107, 61, 0.2) !important;
    outline: none !important;
}

.dark, body {
    background: var(--bg-deep) !important;
}

/* All panels / blocks */
.block, .panel, .form, .tabitem {
    background: var(--bg-card) !important;
    border: 1px solid var(--border-subtle) !important;
    border-radius: 14px !important;
}
.transparent-chatbot, .transparent-chatbot .bubble-wrap, .transparent-chatbot .panel-wrap {
    background: transparent !important;
    border: none !important;
    box-shadow: none !important;
}

/* Tabs */
.tabs > .tab-nav > button {
    font-family: 'Outfit', sans-serif !important;
    font-weight: 600 !important;
    color: var(--text-muted) !important;
    border: none !important;
    background: transparent !important;
    transition: var(--transition) !important;
    padding: 10px 20px !important;
    border-radius: 10px 10px 0 0 !important;
}
.tabs > .tab-nav > button.selected {
    color: var(--accent) !important;
    background: var(--bg-surface) !important;
    border-bottom: 2px solid var(--accent) !important;
}
.tabs > .tab-nav > button:hover:not(.selected) {
    color: var(--flame3) !important;
    background: rgba(255, 107, 61, 0.05) !important;
}

/* Inputs */
input, textarea, .input-text, select {
    background: rgba(15, 12, 20, 0.8) !important;
    border: 1px solid var(--border-subtle) !important;
    color: var(--text-primary) !important;
    border-radius: 10px !important;
    font-family: 'DM Sans', sans-serif !important;
    transition: var(--transition) !important;
}
input:focus, textarea:focus {
    border-color: var(--accent) !important;
    box-shadow: 0 0 0 2px var(--accent-glow) !important;
}

/* Buttons */
button.primary {
    background: linear-gradient(135deg, var(--flame1), var(--flame2), var(--flame3)) !important;
    border: none !important;
    color: white !important;
    font-family: 'Outfit', sans-serif !important;
    font-weight: 700 !important;
    border-radius: 12px !important;
    padding: 10px 24px !important;
    transition: var(--transition) !important;
    box-shadow: 0 4px 15px rgba(255, 69, 0, 0.25) !important;
}
button.primary:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 6px 25px rgba(255, 69, 0, 0.4) !important;
}
button.secondary, button:not(.primary):not(.selected) {
    background: var(--bg-surface) !important;
    border: 1px solid var(--border-subtle) !important;
    color: var(--text-primary) !important;
    border-radius: 12px !important;
    transition: var(--transition) !important;
}
button.secondary:hover, button:not(.primary):not(.selected):hover {
    border-color: var(--accent) !important;
    background: rgba(255, 107, 61, 0.08) !important;
}

/* Chatbot */
.chatbot, .form {
    background: var(--bg-card) !important;
    border: 1px solid var(--border-subtle) !important;
    border-radius: 16px !important;
}
.user-row {
    background: linear-gradient(135deg, rgba(255, 107, 61, 0.15), rgba(255, 69, 0, 0.08)) !important;
    border: 1px solid rgba(255, 107, 61, 0.2) !important;
    border-radius: 16px 16px 4px 16px !important;
}
.bot-row {
    background: var(--bg-surface) !important;
    border: 1px solid var(--border-subtle) !important;
    border-radius: 16px 16px 16px 4px !important;
}
/* Dynamic Text Direction (RTL/LTR Auto-detect) */
.user-row, .bot-row, 
.user-row .prose, .bot-row .prose, 
.user-row .prose *, .bot-row .prose * {
    unicode-bidi: plaintext !important;
    text-align: start !important;
}


@keyframes gradientShift {
    0%   { background-position: 0% 50%; }
    50%  { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
}

@keyframes fireGlow {
    0%, 100% { text-shadow: 0 0 10px rgba(255, 107, 61, 0.5), 0 0 20px rgba(255, 69, 0, 0.3); }
    50% { text-shadow: 0 0 20px rgba(255, 107, 61, 0.8), 0 0 40px rgba(255, 69, 0, 0.5), 0 0 60px rgba(255, 165, 0, 0.3); }
}

@keyframes float {
    0%, 100% { transform: translateY(0); }
    50% { transform: translateY(-12px); }
}

@keyframes fadeInUp {
    from { opacity: 0; transform: translateY(20px); }
    to { opacity: 1; transform: translateY(0); }
}

@keyframes pulse {
    0%, 100% { opacity: 0.6; transform: scale(1); }
    50% { opacity: 1; transform: scale(1.05); }
}

.generating-anim {
    color: var(--gold);
    font-weight: bold;
    animation: pulse 1.5s infinite;
    padding: 10px;
}

@keyframes ember {
    0% { transform: translateY(0) scale(1); opacity: 0.8; }
    50% { transform: translateY(-30px) scale(0.8); opacity: 0.4; }
    100% { transform: translateY(-60px) scale(0.5); opacity: 0; }
}

/* ── Chat Placeholder ──────────────────────────────────── */
.placeholder-content {
    display: flex !important;
    flex-direction: column !important;
    align-items: center !important;
    justify-content: center !important;
    text-align: center !important;
    animation: fadeInUp 0.8s ease-out;
}
.placeholder-content h1 {
    font-family: 'Outfit', sans-serif;
    font-size: 2.2rem; font-weight: 800;
    background: linear-gradient(135deg, #ffd700, #ff6b3d, #ff4500);
    background-size: 200% 200%;
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    animation: gradientShift 4s ease infinite;
    margin: 0 0 8px 0;
    display: flex;
    flex-direction: column;
    align-items: center;
}
.placeholder-content h2 {
    color: var(--text-muted); font-size: 1rem;
    font-family: 'DM Sans', sans-serif;
    letter-spacing: 0.03em;
    font-weight: normal;
    margin-bottom: 24px;
}
.placeholder-content p {
    background: rgba(255, 107, 61, 0.1);
    border: 1px solid rgba(255, 107, 61, 0.2);
    padding: 10px 20px;
    border-radius: 999px;
    display: inline-block;
    color: var(--accent);
    font-weight: 600;
    font-size: 0.95rem;
    letter-spacing: 0.02em;
}

/* ── Password status ──────────────────────────────────── */
.pw-status-box {
    margin-top: 8px; padding: 10px 16px; border-radius: 12px;
    font-family: 'Outfit', sans-serif; font-size: 0.88rem; font-weight: 700;
    text-align: center; transition: var(--transition);
}
.pw-locked { background: rgba(248,113,113,0.1); color: #fca5a5; border: 1px solid rgba(248,113,113,0.3); }
.pw-ok { background: rgba(34,212,122,0.1); color: #86efac; border: 1px solid rgba(34,212,122,0.3); }

/* ── Header ───────────────────────────────────────────── */
.ph-header.flame {
    background: linear-gradient(135deg, rgba(20, 8, 5, 0.97) 0%, rgba(15, 8, 18, 0.97) 40%, rgba(10, 10, 14, 0.98) 100%);
    border: 1px solid rgba(255, 110, 60, 0.15);
    box-shadow: 0 12px 50px rgba(255, 69, 0, 0.08), 0 4px 20px rgba(0,0,0,0.4);
    padding: 28px 34px;
    border-radius: 18px;
    margin: 14px 14px 18px 14px;
    position: relative;
    overflow: hidden;
}
.ph-header.flame::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 3px;
    background: linear-gradient(90deg, var(--flame1), var(--flame3), var(--flame4), var(--flame2));
    background-size: 300% 100%;
    animation: gradientShift 4s ease infinite;
}
.ph-header.flame::after {
    content: '';
    position: absolute;
    bottom: 0; left: 50%; transform: translateX(-50%);
    width: 80%; height: 1px;
    background: linear-gradient(90deg, transparent, rgba(255, 107, 61, 0.3), transparent);
}
.ph-title {
    font-family: 'Outfit', sans-serif;
    font-size: 2.4rem;
    font-weight: 800;
    letter-spacing: -0.03em;
    line-height: 1.05;
    background: linear-gradient(90deg, #ffd700 0%, #ff7a3d 25%, #ff4500 50%, #ff7a3d 75%, #ffd700 100%);
    background-size: 300% 100%;
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    animation: gradientShift 5s ease infinite;
}
.ph-version { font-size: 0.45em; opacity: 0.5; font-weight: 600; margin-left: 6px; }
.ph-sub {
    color: var(--text-muted); font-size: 0.88rem; margin-top: 8px;
    font-family: 'DM Sans', sans-serif; opacity: 0.85;
    letter-spacing: 0.02em;
}
.ph-pills { margin-top: 14px; display: flex; gap: 8px; flex-wrap: wrap; }
.pill {
    display: inline-flex; align-items: center; gap: 5px; padding: 5px 13px;
    border-radius: 999px; font-size: 0.72rem; font-weight: 700;
    font-family: 'Outfit', sans-serif; backdrop-filter: blur(8px);
    transition: var(--transition);
}
.pill:hover { transform: translateY(-1px); }
.pill-red    { background: rgba(255,69,0,0.1); color: #ffb2a0; border: 1px solid rgba(255,69,0,0.2); }
.pill-blue   { background: rgba(108,114,255,0.1); color: #c5cbff; border: 1px solid rgba(108,114,255,0.2); }
.pill-green  { background: rgba(34,212,122,0.1); color: #b5f5cf; border: 1px solid rgba(34,212,122,0.2); }
.pill-gold   { background: rgba(255,165,0,0.1); color: #ffe6b0; border: 1px solid rgba(255,165,0,0.2); }

/* ── Examples grid ─────────────────────────────────────── */
.gradio-container .examples button,
.gradio-container button.example,
.gradio-container .gallery-item {
    background: linear-gradient(135deg, rgba(255, 69, 0, 0.05), rgba(255, 165, 0, 0.03)) !important;
    border: 1px solid rgba(255, 107, 61, 0.15) !important;
    border-radius: 12px !important;
    color: var(--text-primary) !important;
    font-size: 0.85rem !important;
    font-family: 'DM Sans', sans-serif !important;
    padding: 12px 18px !important;
    transition: var(--transition) !important;
    text-align: left !important;
    backdrop-filter: blur(4px) !important;
}
.gradio-container .examples button:hover,
.gradio-container button.example:hover,
.gradio-container .gallery-item:hover {
    border-color: var(--accent) !important;
    background: linear-gradient(135deg, rgba(255, 69, 0, 0.12), rgba(255, 165, 0, 0.08)) !important;
    transform: translateY(-3px) !important;
    box-shadow: 0 8px 25px rgba(255, 69, 0, 0.15) !important;
}



/* ── Chat Placeholder ──────────────────────────────────── */
.placeholder-content {
    display: flex !important;
    flex-direction: column !important;
    align-items: center !important;
    justify-content: center !important;
    height: 100% !important;
    text-align: center !important;
    background: transparent !important;
    border: none !important;
    box-shadow: none !important;
    animation: fadeInUp 0.8s ease-out;
    position: relative;
    padding: 40px 20px;
}
.placeholder-content .prose {
    background: transparent !important;
    border: none !important;
    box-shadow: none !important;
}
.placeholder-content h1::before {
    content: "🦅";
    font-size: 4.5rem;
    animation: float 3s ease-in-out infinite;
    background: radial-gradient(circle, rgba(255, 107, 61, 0.45) 0%, transparent 60%);
    display: inline-block;
    width: 130px;
    height: 130px;
    line-height: 130px;
    margin-bottom: 16px;
    border-radius: 50%;
    -webkit-text-fill-color: initial !important;
}
.placeholder-content h1 {
    font-family: 'Outfit', sans-serif !important;
    font-size: 2.2rem !important;
    font-weight: 800 !important;
    background: linear-gradient(135deg, #ffd700, #ff6b3d, #ff4500) !important;
    background-size: 200% 200% !important;
    -webkit-background-clip: text !important;
    -webkit-text-fill-color: transparent !important;
    animation: gradientShift 4s ease infinite;
    margin: 0 0 8px 0 !important;
}
.placeholder-content h2 {
    color: var(--text-muted) !important;
    font-size: 1rem !important;
    font-family: 'DM Sans', sans-serif !important;
    letter-spacing: 0.03em !important;
    margin-bottom: 24px !important;
    font-weight: normal !important;
}
.placeholder-content p {
    background: rgba(255, 107, 61, 0.1) !important;
    border: 1px solid rgba(255, 107, 61, 0.2) !important;
    padding: 10px 20px !important;
    border-radius: 999px !important;
    display: inline-block !important;
    color: var(--text-primary) !important;
    font-size: 0.85rem !important;
    font-family: 'DM Sans', sans-serif !important;
    animation: fadeInUp 1s ease-out;
    margin: 0 !important;
}
::-webkit-scrollbar { width: 6px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb {
    background: rgba(255, 107, 61, 0.2);
    border-radius: 3px;
}
::-webkit-scrollbar-thumb:hover { background: rgba(255, 107, 61, 0.4); }

/* ── Labels ────────────────────────────────────────────── */
label, .label-wrap {
    color: var(--text-muted) !important;
    font-family: 'DM Sans', sans-serif !important;
    font-size: 0.85rem !important;
}

/* ── File upload area ──────────────────────────────────── */
.upload-container, .file-upload {
    border: 2px dashed rgba(255, 107, 61, 0.15) !important;
    border-radius: 14px !important;
    background: rgba(15, 12, 20, 0.5) !important;
    transition: var(--transition) !important;
}
.upload-container:hover, .file-upload:hover {
    border-color: var(--accent) !important;
    background: rgba(255, 107, 61, 0.03) !important;
}

/* ── Media tab ─────────────────────────────────────────── */
.audio-player, .video-player, audio, video {
    border-radius: 12px !important;
    border: 1px solid var(--border-subtle) !important;
}

@media (max-width: 768px) {
    .ph-title { font-size: 1.6rem; }
    .pill { font-size: 0.65rem; padding: 4px 8px; }
    .placeholder-title { font-size: 1.6rem; }
    .placeholder-features { gap: 8px; }
}
"""

# ==============================================================================
# 🏷️  HEADER HTML + SIDEBAR GUIDE
# ==============================================================================
# Defines the raw HTML elements and user instruction manual for the application interface.

HEADER_HTML = """
<div class="ph-header flame">
  <div class="ph-title">Phoenix Eduplan <span class="ph-version">v4.1</span></div>
  <div class="ph-sub">Advanced AI study assistant · Persistent storage · Multi-doc RAG · Hybrid retrieval</div>
  <div class="ph-pills">
    <span class="pill pill-red">🦅 Phoenix Core</span>
    <span class="pill pill-blue">⚡ Llama 3.3 70B</span>
    <span class="pill pill-green">🔍 Hybrid Search</span>
    <span class="pill pill-gold">📄 PDFs · DOCX · PPTX · XLSX · CSV · Images</span>
  </div>
</div>
"""

SIDEBAR_GUIDE = """
### 💡 Quick Commands
| Say... | Gets you... |
|--------|------------|
| `summarize` | Structured summary |
| `quiz me` | 5 MCQ questions |
| `flashcards` | 6 study cards |
| `compare X and Y` | Side-by-side comparison |
| `mind map` | Hierarchical concept map |
| `explain simply` | ELI15 explanation |

### 🎨 New in v4.1
- 🔧 Fixed multi-upload document persistence
- 📊 Larger context window (12K chars)
- 🎯 Improved retrieval with better chunking (800/200)
- 🔍 Better deduplication in hybrid search
- 🎨 Glassmorphism UI with micro-animations
- ⚡ Faster streaming with batched output
- 🗄️ More robust database connections
- 📝 Query result caching
"""

# ==============================================================================
# 🖥️  BUILD UI
# ==============================================================================
# The massive layout orchestrator that links all Python backend logic to the Gradio frontend event loops.

def build_ui():
    js_func = """
    function refresh() {
        const url = new URL(window.location);
        if (url.searchParams.get('__theme') !== 'dark') {
            url.searchParams.set('__theme', 'dark');
            window.location.href = url.href;
        }
    }
    """
    with gr.Blocks(
        title="Phoenix Eduplan v4.1",
        css=DARK_CSS,
        theme=gr.themes.Default(primary_hue="orange", neutral_hue="zinc"),
        js=js_func
    ) as demo:
        user_state = gr.State("")

        with gr.Column(visible=True, elem_classes=["login-container"]) as login_view:
            gr.HTML("""
                <div style="margin-bottom: 1rem;">
                    <span style="font-size: 3rem;">🦅</span>
                </div>
                <div class="login-header-text">Phoenix Eduplan</div>
                <div class="login-subtitle">Enter your secure study environment</div>
            """)
            login_user = gr.Textbox(label="Username", placeholder="Enter your username", elem_classes=["custom-input"])
            login_pass = gr.Textbox(label="Password (Required)", placeholder="Enter password", type="password", elem_classes=["custom-input"])
            login_btn = gr.Button("Sign In to Continue", variant="primary", size="lg")

        with gr.Column(visible=False) as main_view:

            # ── Header ───────────────────────────────────────────────────────────
            gr.HTML(HEADER_HTML)

            with gr.Row():
                # ── Sidebar (left column) ────────────────────────────────────────
                with gr.Column(scale=1, min_width=280):
                    gr.Markdown("### 🔐 Access")
                    logout_btn = gr.Button("Logout (Switch User)")

                    gr.Markdown("---")

                    with gr.Tabs():
                        with gr.Tab("📁 Documents"):
                            file_input = gr.File(
                                label="Upload Files",
                                file_count="multiple",
                                file_types=[
                                    ".pdf", ".docx", ".doc", ".pptx", ".ppt",
                                    ".xlsx", ".xls", ".csv", ".txt", ".md",
                                    ".png", ".jpg", ".jpeg", ".json", ".xml",
                                ],
                                interactive=False,
                            )
                            process_btn = gr.Button("📥 Process Files", variant="primary", interactive=False)
                            clear_btn   = gr.Button("🗑️ Clear All", interactive=False)
                            doc_status  = gr.Markdown("📁 No documents loaded.")

                        with gr.Tab("📋 Sessions"):
                            session_list    = gr.Dropdown(label="Switch Session", choices=[], interactive=True)
                            new_session_btn = gr.Button("➕ New Session")
                            session_info    = gr.Markdown("")

                        with gr.Tab("💾 Data"):
                            download_db_btn   = gr.Button("📥 Download DB", interactive=False)
                            download_hist_btn = gr.Button("📥 Download History", interactive=False)
                            db_file           = gr.File(label="Database")
                            hist_file         = gr.File(label="History ZIP")

                        with gr.Tab("ℹ️ Guide"):
                            gr.Markdown(SIDEBAR_GUIDE)

                # ── Main area (right column) ─────────────────────────────────────
                with gr.Column(scale=3):
                    with gr.Tabs():
                        with gr.Tab("💬 Chat"):
                            chat = gr.ChatInterface(
                                fn=chat_logic,
                                additional_inputs=[user_state],
                                chatbot=gr.Chatbot(
                                    height="70vh",
                                    elem_classes=["transparent-chatbot"],
                                    placeholder='''
<div style="text-align: center;">
    <h1>Phoenix Eduplan v4.1</h1>
    <h3>Upload your documents and ignite your learning</h3>
    <p>📚 Multi-doc RAG  •  ⚡ AI-Powered  •  🌍 Bilingual</p>
</div>
''',
                                ),
                                examples=[
                                    ["Deep dive into this topic"],
                                    ["تعمق في هذا الموضوع"],
                                    ["Teach me step by step"],
                                    ["علمني خطوة بخطوة"],
                                    ["Generate 10 practice problems"],
                                    ["أعطني 10 مسائل للتدريب"],
                                    ["Create a comprehensive study guide"],
                                    ["أنشئ دليل دراسي شامل"],
                                    ["Debate the pros and cons"],
                                    ["ناقش الإيجابيات والسلبيات"],
                                ],
                                cache_examples=False,
                            )

                        with gr.Tab("🎧 Media"):
                            gr.Markdown("### 🎧 Audio & Video Generation")
                            media_text = gr.Textbox(
                                label="Script / Text",
                                placeholder="Enter text for audio/video, or leave empty to use uploaded documents...",
                                lines=5,
                                interactive=False,
                            )
                            lang_radio = gr.Radio(["auto", "en", "ar"], label="Language", value="auto", interactive=False)
                            with gr.Tabs():
                                with gr.Tab("🎵 Audio"):
                                    audio_btn = gr.Button("🔊 Generate Audio", variant="primary", interactive=False)
                                    audio_output = gr.Audio(label="Generated Audio", interactive=False)
                                with gr.Tab("🎥 Video"):
                                    video_btn = gr.Button("🎬 Generate AI Video", variant="primary", interactive=False)
                                    video_output = gr.Video(label="Generated Video", interactive=False)
                        
                        with gr.Tab("📋 Study Plan"):
                            gr.Markdown("### 📋 AI Study Plan Generator")
                            plan_duration = gr.Radio(["1 week", "2 weeks", "1 month"], label="Duration", value="2 weeks", interactive=False)
                            plan_lang = gr.Radio(["auto", "en", "ar"], label="Language", value="auto", interactive=False)
                            plan_btn = gr.Button("🔄 Generate Study Plan", variant="primary", interactive=False)
                            plan_output = gr.Markdown("Click generate to create a plan from your active documents.")
                        
                        with gr.Tab("✍️ Essay Grader"):
                            gr.Markdown("### ✍️ AI Essay Grader")
                            essay_text = gr.Textbox(label="Paste your essay here", lines=10, interactive=False)
                            rubric_text = gr.Textbox(label="Optional Grading Rubric", lines=2, interactive=False)
                            essay_lang = gr.Radio(["auto", "en", "ar"], label="Language", value="auto", interactive=False)
                            grade_btn = gr.Button("✅ Grade Essay", variant="primary", interactive=False)
                            grade_output = gr.Markdown("")
                        
                        with gr.Tab("📊 Analytics & Vocab"):
                            gr.Markdown("### 📊 Learning Analytics")
                            analytics_btn = gr.Button("🔄 Refresh Analytics", interactive=False)
                            analytics_output = gr.Markdown("")
                            gr.Markdown("---")
                            gr.Markdown("### 📚 Vocabulary Extractor")
                            vocab_lang = gr.Radio(["auto", "en", "ar"], label="Language", value="auto", interactive=False)
                            vocab_btn = gr.Button("📝 Extract Vocabulary", interactive=False)
                            vocab_output = gr.Markdown("")

            # ── Access control outputs ───────────────────────────────────────────
            _shared_access_outputs = [
                file_input, process_btn, clear_btn,
                download_db_btn, download_hist_btn,
                media_text, lang_radio, audio_btn, video_btn,
                plan_duration, plan_lang, plan_btn,
                essay_text, rubric_text, essay_lang, grade_btn,
                analytics_btn, vocab_lang, vocab_btn
            ]

            def on_page_load(username):
                if not username:
                    return (
                        *[gr.update(interactive=False) for _ in range(len(_shared_access_outputs))],
                        gr.update(choices=[]), "📁 Please log in.", None
                    )
                ensure_user_dirs(username)
                sid = get_or_create_user_session(username)
                sessions = list_user_sessions(username)
                choices = [f"{s['name']} ({s['id']})" for s in sessions]
                doc_text = list_documents(username)
                history = load_user_history(username)
                return (
                    *[gr.update(interactive=True) for _ in range(len(_shared_access_outputs))],
                    gr.update(choices=choices, value=f"{sessions[0]['name']} ({sid})" if sessions else None), 
                    doc_text, history if history else None
                )


            user_state.change(
                fn=on_page_load,
                inputs=[user_state],
                outputs=[
                    *_shared_access_outputs,
                    session_list, doc_status, chat.chatbot,
                ],
            )
            demo.load(
                fn=on_page_load,
                inputs=[user_state],
                outputs=[
                    *_shared_access_outputs,
                    session_list, doc_status, chat.chatbot,
                ],
            )

            # ── File processing ──────────────────────────────────────────────────
            def on_process(files, username):
                if not username: return "⚠️ Please log in."
                result = process_files(files, username)
                gr.Info(result)
                doc_text = list_documents(username)
                return doc_text

            def start_processing():
                return gr.update(value="⏳ Processing...", interactive=False)
            
            def end_processing():
                return gr.update(value="📥 Process Files", interactive=True)

            process_btn.click(
                fn=start_processing, inputs=[], outputs=[process_btn]
            ).then(
                fn=on_process,
                inputs=[file_input, user_state],
                outputs=[doc_status],
            ).then(
                fn=end_processing, inputs=[], outputs=[process_btn]
            )

            # ── Clear documents ──────────────────────────────────────────────────
            def on_clear(username):
                if not username: return "⚠️ Please log in.", gr.update(), gr.update(), None
                msg, _ = clear_documents(username)
                
                # Force a new session when clearing documents so the LLM forgets old chats about them
                sid = create_session(username=username)
                ctx = get_user_context(username)
                ctx["session_id"] = sid
                sessions = list_user_sessions(username)
                choices = [f"{s['name']} ({s['id']})" for s in sessions]
                
                return msg, gr.update(choices=choices), f"✅ New session: {sid}", None

            clear_btn.click(
                fn=on_clear,
                inputs=[user_state],
                outputs=[doc_status, session_list, session_info, chat.chatbot],
            )

            # ── Session management ───────────────────────────────────────────────
            def on_new_session(username):
                if not username:
                    return gr.update(), "⚠️ Log in first.", None
                sid = create_session(username=username)
                ctx = get_user_context(username)
                ctx["session_id"] = sid
                sessions = list_user_sessions(username)
                choices = [f"{s['name']} ({s['id']})" for s in sessions]
                return gr.update(choices=choices), f"✅ New session: {sid}", None

            new_session_btn.click(
                fn=on_new_session,
                inputs=[user_state],
                outputs=[session_list, session_info, chat.chatbot],
            )

            def on_switch_session(selection, username):
                if not selection or not username:
                    return "⚠️ Select a valid session.", None
                match = re.search(r"\(([^)]+)\)$", selection)
                if not match:
                    return "⚠️ Invalid selection.", None
                sid = match.group(1)
                history = switch_user_session(username, sid)
                return f"✅ Switched to session {sid}", history if history else None

            session_list.change(
                fn=on_switch_session,
                inputs=[session_list, user_state],
                outputs=[session_info, chat.chatbot],
            )

            # ── Downloads ────────────────────────────────────────────────────────
            def on_download_db():
                path = download_db()
                if path:
                    return gr.update(value=path)
                return gr.update()

            download_db_btn.click(fn=on_download_db, outputs=[db_file])

            def on_download_hist(username):
                if not username: return gr.update()
                path = download_user_history(username)
                if path:
                    return gr.update(value=path)
                return gr.update()

            download_hist_btn.click(
                fn=on_download_hist,
                inputs=[user_state],
                outputs=[hist_file],
            )

            # ── Media generation ─────────────────────────────────────────────────
            def start_audio(): return gr.update(value="⏳ Generating Audio...", interactive=False)
            def end_audio(audio_path): 
                return gr.update(value="🔊 Generate Audio", interactive=True)
            def on_gen_audio(text, lang, username):
                if not username: return gr.update(value=None)
                audio_path, msg = generate_audio_only(text, username, lang)
                if audio_path:
                    gr.Info(msg)
                else:
                    gr.Warning(msg)
                return gr.update(value=audio_path)

            audio_btn.click(
                fn=start_audio, inputs=[], outputs=[audio_btn]
            ).then(
                fn=on_gen_audio,
                inputs=[media_text, lang_radio, user_state],
                outputs=[audio_output],
            ).then(
                fn=end_audio, inputs=[audio_output], outputs=[audio_btn]
            )

            def start_video(): return gr.update(value="🎬 Generating Video (This may take a while)...", interactive=False)
            def end_video(video_path):
                return gr.update(value="🎬 Generate AI Video", interactive=True)
            def on_gen_video(text, lang, username):
                if not username: return gr.update(value=None)
                video_path, msg = generate_video_only(text, username, lang)
                if video_path:
                    gr.Info(msg)
                else:
                    gr.Warning(msg)
                return gr.update(value=video_path)

            video_btn.click(
                fn=start_video, inputs=[], outputs=[video_btn]
            ).then(
                fn=on_gen_video,
                inputs=[media_text, lang_radio, user_state],
                outputs=[video_output],
            ).then(
                fn=end_video, inputs=[video_output], outputs=[video_btn]
            )
        
            # ── New v5.0 Features ────────────────────────────────────────────────
        
            def start_plan(): return gr.update(value="⏳ Generating Plan...", interactive=False)
            def end_plan(): return gr.update(value="🔄 Generate Study Plan", interactive=True)
            def on_gen_plan(duration, lang, username):
                if username:
                    return generate_study_plan(username, duration, lang)
                return "⚠️ Please log in."

            plan_btn.click(fn=start_plan, inputs=[], outputs=[plan_btn]).then(
                fn=on_gen_plan,
                inputs=[plan_duration, plan_lang, user_state],
                outputs=[plan_output]
            ).then(fn=end_plan, inputs=[], outputs=[plan_btn])
        
            def start_grade(): return gr.update(value="⏳ Grading...", interactive=False)
            def end_grade(): return gr.update(value="✅ Grade Essay", interactive=True)
            grade_btn.click(fn=start_grade, inputs=[], outputs=[grade_btn]).then(
                fn=grade_essay,
                inputs=[essay_text, rubric_text, essay_lang],
                outputs=[grade_output]
            ).then(fn=end_grade, inputs=[], outputs=[grade_btn])
        
            def on_analytics(username):
                if not username: return "⚠️ Please log in."
                return get_analytics(username)
            analytics_btn.click(
                fn=on_analytics,
                inputs=[user_state],
                outputs=[analytics_output]
            )
        
            def start_vocab(): return gr.update(value="⏳ Extracting...", interactive=False)
            def end_vocab(): return gr.update(value="📝 Extract Vocabulary", interactive=True)
            def on_gen_vocab(lang, username):
                if username:
                    return extract_vocabulary(username, lang)
                return "⚠️ Please log in."

            vocab_btn.click(fn=start_vocab, inputs=[], outputs=[vocab_btn]).then(
                fn=on_gen_vocab,
                inputs=[vocab_lang, user_state],
                outputs=[vocab_output]
            ).then(fn=end_vocab, inputs=[], outputs=[vocab_btn])

    
        # ── Auth Event Logic ───────────────────────────────────────────────
        def handle_login(user, pwd):
            """Authenticates users using a memory-hard 'scrypt' Key Derivation Function. Uses constant-time string comparison to mitigate cryptographic timing attacks."""
            import secrets
            if not user or len(user.strip()) < 3:
                return gr.update(), gr.update(), gr.update(), "⚠️ Username must be at least 3 characters."
            
            # Use a memory-hard scrypt KDF to make offline brute-forcing impossible
            expected_hash = b'\xcdf\x20\xde\x0b\xc8\xc2\xcb\x8a\xcc\xa7\xfb\xff\x91\x8d\xa3\x91\x1d\x96\x1e\x73\xe6\x14\x75\x14\xad\x72\xa0\"\xf9\xbe\xff'
            try:
                actual_hash = hashlib.scrypt(pwd.encode('utf-8'), salt=b"phoenix_eduplan_salt_2026", n=16384, r=8, p=1, dklen=32)
                is_valid = secrets.compare_digest(actual_hash, expected_hash)
            except Exception:
                is_valid = False

            if not is_valid:
                return gr.update(), gr.update(), gr.update(), "⚠️ Incorrect password."
            return user.strip(), gr.update(visible=False), gr.update(visible=True), ""
            
        def handle_logout():
            return "", gr.update(visible=True), gr.update(visible=False)

        login_btn.click(
            fn=handle_login,
            inputs=[login_user, login_pass],
            outputs=[user_state, login_view, main_view, login_user]
        )
        logout_btn.click(
            fn=handle_logout,
            inputs=[],
            outputs=[user_state, login_view, main_view]
        )

        # ── Dedicated Backend APIs ───────────────────────────────────────────
        django_username = gr.Textbox(visible=False)
        django_dummy_btn = gr.Button(visible=False)
        
        django_dummy_btn.click(
            fn=lambda t, l, u: generate_audio_only(t, u, l)[0],
            inputs=[media_text, lang_radio, django_username],
            outputs=[audio_output],
            api_name="django_audio"
        )
        django_dummy_btn.click(
            fn=lambda t, l, u: generate_video_only(t, u, l)[0],
            inputs=[media_text, lang_radio, django_username],
            outputs=[video_output],
            api_name="django_video"
        )
        django_dummy_btn.click(
            fn=lambda d, l, u: generate_study_plan(u, d, l),
            inputs=[plan_duration, plan_lang, django_username],
            outputs=[plan_output],
            api_name="django_plan"
        )
        django_dummy_btn.click(
            fn=lambda u: get_analytics(u),
            inputs=[django_username],
            outputs=[analytics_output],
            api_name="django_analytics"
        )
        django_dummy_btn.click(
            fn=lambda l, u: extract_vocabulary(u, l),
            inputs=[vocab_lang, django_username],
            outputs=[vocab_output],
            api_name="django_vocab"
        )

    return demo

# ==============================================================================
# 🚀  LAUNCH + API
# ==============================================================================
# Bootstraps the application and starts the web server.

demo = build_ui()

if __name__ == "__main__":
    demo.launch(
        server_name="0.0.0.0",
        server_port=7860,
        share=False,
        show_error=True,
    )